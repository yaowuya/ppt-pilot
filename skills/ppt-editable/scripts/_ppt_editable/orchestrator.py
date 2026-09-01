"""Fixed-order orchestration for editable PowerPoint generation and promotion."""

from __future__ import annotations

from dataclasses import asdict, dataclass, replace
import importlib.util
import os
from pathlib import Path
import shutil
import stat
import sys
from typing import Callable, Mapping, Optional, Tuple
import uuid

from .atomic_io import (
    OutputLock,
    _manifest_files_match,
    _quarantine_file,
    _read_json,
    _result_manifest_is_valid,
    atomic_write_bytes,
    begin_promotion,
    build_output_paths,
    promote_output,
    recover_incomplete_transactions,
)
from .config import load_verification_config
from .contract import (
    RunContext,
    SlideSource,
    locate_run,
    parse_storyboard,
    resolve_slide_sources,
    validate_completed_run,
    validate_safe_regular_file,
)
from .errors import EditableError
from .model import DeckPlan, EditableResult, Failure
from .office_protocol import (
    OfficeResult,
    RENDER_KEYS,
    invoke_office_verification,
)
from .snapshot import canonical_snapshot_payload, compute_snapshot_id


CONVERTER_VERSION = "ppt-editable-converter-v1"
SUBSET_CONTRACT_VERSION = "editable-svg-subset-v1"


@dataclass(frozen=True)
class GenerationCapability:
    office_available: bool
    pillow_available: bool
    office_runner: Optional[Callable[[Mapping[str, object]], OfficeResult]] = None

    def __post_init__(self) -> None:
        if type(self.office_available) is not bool or type(self.pillow_available) is not bool:
            raise TypeError("capability flags must be boolean")
        if self.office_runner is not None and not callable(self.office_runner):
            raise TypeError("office_runner must be callable")


def _failure(
    code: str,
    message: str,
    *,
    slide_id: Optional[str] = None,
) -> Failure:
    return Failure(
        code=code,
        slide_id=slide_id,
        svg_tree_path=None,
        element_type=None,
        message=message,
        remediation="repair the completed run or verification capability and retry",
    )


def _error_failure(error: EditableError) -> Failure:
    return Failure(
        code=error.code,
        slide_id=error.slide_id,
        svg_tree_path=error.svg_tree_path,
        element_type=error.element_type,
        message=error.message,
        remediation=error.remediation or "repair the completed run and retry",
    )


def _result(
    status: str,
    deck_id: str,
    snapshot_id: str,
    slide_count: int,
    *,
    failures: Tuple[Failure, ...] = (),
    warnings: Tuple[str, ...] = (),
) -> EditableResult:
    return EditableResult(
        status=status,
        deck_id=deck_id,
        input_snapshot_id=snapshot_id,
        slide_count=slide_count,
        failures=failures,
        warnings=warnings,
    )


def editable_result_dict(result: EditableResult) -> Mapping[str, object]:
    value = asdict(result)
    value["schema_version"] = 1
    value["kind"] = "ppt_editable_invocation"
    return value


def _committed_result(paths, snapshot_id: str) -> Optional[EditableResult]:
    exists = os.path.lexists(str(paths.manifest_path))
    manifest = _read_json(paths.manifest_path)
    if not exists:
        return None
    if (
        manifest is None
        or not _result_manifest_is_valid(manifest, paths)
        or not _manifest_files_match(manifest, paths)
    ):
        raise EditableError(
            "promotion_conflict",
            "editable-result.json does not describe coherent public authority",
        )
    if manifest.get("input_snapshot_id") != snapshot_id:
        return None
    return EditableResult(
        status=str(manifest["status"]),
        deck_id=str(manifest["deck_id"]),
        input_snapshot_id=str(manifest["input_snapshot_id"]),
        slide_count=int(manifest["slide_count"]),
        output_path=str(manifest["output_path"]),
        output_sha256=str(manifest["output_sha256"]),
        failures=(),
        warnings=tuple(manifest.get("warnings", ())),
    )


def _build_presentation_bytes(deck_plan: DeckPlan, *, include_text: bool = True) -> bytes:
    from .drawingml import presentation_bytes

    return presentation_bytes(deck_plan, include_text=include_text)


def _verify_candidate(
    path: Path,
    deck_plan: DeckPlan,
    config,
    *,
    office_normalized: bool = False,
):
    from .structural_verify import verify_candidate

    return verify_candidate(
        path,
        deck_plan,
        config,
        office_normalized=office_normalized,
    )


def _verify_geometry_candidate(
    path: Path,
    deck_plan: DeckPlan,
) -> Tuple[Failure, ...]:
    import zipfile
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    failures = []
    try:
        with zipfile.ZipFile(str(path), "r") as archive:
            if archive.testzip() is not None:
                failures.append(_failure("pptx_zip_invalid", "geometry candidate ZIP is corrupt"))
            if any(name.startswith("ppt/media/") for name in archive.namelist()):
                failures.append(_failure("image_fallback_detected", "geometry candidate contains media"))
        presentation = Presentation(str(path))
    except Exception as exc:
        return (_failure("pptx_reopen_failed", "geometry candidate cannot reopen: {}".format(exc)),)
    if len(presentation.slides) != len(deck_plan.slides):
        failures.append(_failure("structure_mismatch", "geometry candidate slide count differs"))
    if (presentation.slide_width, presentation.slide_height) != (12192000, 6858000):
        failures.append(_failure("structure_mismatch", "geometry candidate dimensions differ"))

    def inspect_shape(shape) -> None:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
            failures.append(_failure("content_mismatch", "geometry candidate contains visible text"))
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            failures.append(_failure("image_fallback_detected", "geometry candidate contains picture"))
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if len(shape.shapes) == 0:
                failures.append(_failure("group_mismatch", "geometry candidate contains empty group"))
            for child in shape.shapes:
                inspect_shape(child)

    for slide in presentation.slides:
        for shape in slide.shapes:
            inspect_shape(shape)
    return tuple(failures)


def _write_candidate_bytes(path: Path, data: bytes) -> str:
    return atomic_write_bytes(path, data)


def _core_dependency_failure() -> Optional[Failure]:
    if sys.version_info < (3, 9):
        return _failure(
            "python_version_unsupported",
            "ppt-editable requires Python 3.9 or newer",
        )
    missing = [
        name
        for name in ("pptx", "defusedxml")
        if importlib.util.find_spec(name) is None
    ]
    if missing:
        return _failure(
            "core_dependency_missing",
            "missing core dependency: {}".format(", ".join(sorted(missing))),
        )
    return None


def _create_work_directory(paths) -> Path:
    directory = paths.tmp_dir / ("work-" + uuid.uuid4().hex)
    os.mkdir(str(directory))
    metadata = os.lstat(str(directory))
    if not stat.S_ISDIR(metadata.st_mode) or getattr(metadata, "st_file_attributes", 0) & 0x400:
        raise EditableError("source_path_unsafe", "created work directory is unsafe")
    return directory


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _geometry_svg_bytes(data: bytes) -> bytes:
    from defusedxml import ElementTree as DET

    root = DET.fromstring(data)

    def prune(parent) -> None:
        for child in list(parent):
            name = _local_name(child.tag) if isinstance(child.tag, str) else ""
            if name == "text":
                parent.remove(child)
                continue
            prune(child)
            if name == "g" and not list(child):
                parent.remove(child)

    prune(root)
    return DET.tostring(root, encoding="utf-8", xml_declaration=True)


def _write_geometry_svgs(
    sources: Tuple[SlideSource, ...],
    directory: Path,
) -> Tuple[Mapping[str, str], ...]:
    directory.mkdir(parents=True, exist_ok=False)
    result = []
    for source in sources:
        output = directory / (source.slide_id + ".svg")
        _write_candidate_bytes(
            output,
            _geometry_svg_bytes(Path(source.path).read_bytes()),
        )
        result.append({"slide_id": source.slide_id, "path": str(output)})
    return tuple(result)


def _office_failure(result: OfficeResult) -> Failure:
    stage = "" if result.error is None else result.error.get("stage", "")
    if stage == "normalize":
        code = "powerpoint_normalize_failed"
    elif stage in ("counts", "reopen"):
        code = "powerpoint_reopen_failed"
    else:
        code = "powerpoint_render_failed"
    message = (
        "PowerPoint verification failed"
        if result.error is None
        else result.error.get("message", "PowerPoint verification failed")
    )
    return _failure(code, message)


def _validate_office_result_for_request(
    office: OfficeResult,
    request: Mapping[str, object],
) -> Optional[Failure]:
    if not isinstance(office, OfficeResult):
        return _failure(
            "powerpoint_render_failed",
            "Office runner returned an invalid result type",
        )
    if not office.capability or office.exit_code != 0 or office.error is not None:
        return None
    if dict(office.counts) != dict(request["expected_counts"]):
        return _failure(
            "powerpoint_reopen_failed",
            "Office recursive counts differ from request",
        )
    if office.normalized_path != request["normalized_path"]:
        return _failure(
            "powerpoint_normalize_failed",
            "Office normalized path differs from request",
        )
    expected_ids = tuple(request["ordered_slide_ids"])
    for key in RENDER_KEYS:
        entries = tuple(office.renders.get(key, ()))
        if tuple(entry.get("slide_id") for entry in entries) != expected_ids:
            return _failure(
                "powerpoint_render_failed",
                "Office render order differs for {}".format(key),
            )
        directory = Path(str(request["render_directories"][key])).absolute()
        for entry, slide_id in zip(entries, expected_ids):
            if Path(str(entry.get("path", ""))).absolute() != directory / (slide_id + ".png"):
                return _failure(
                    "powerpoint_render_failed",
                    "Office render path differs for {}".format(key),
                )
    expected_stages = (
        ("capability", "passed"),
        ("normalize", "running"),
        ("normalize", "passed"),
        ("counts", "running"),
        ("counts", "passed"),
        ("source_decks", "running"),
        ("source_decks", "passed"),
        ("render", "running"),
        ("render", "passed"),
    )
    actual_stages = tuple(
        (stage.get("name"), stage.get("status"))
        if isinstance(stage, Mapping)
        else (None, None)
        for stage in office.stages
    )
    if actual_stages != expected_stages:
        return _failure(
            "powerpoint_render_failed",
            "Office stage sequence differs from request contract",
        )
    return None


def _quarantine_work(paths, directory: Path, label: str) -> None:
    if not directory.exists():
        return
    destination = paths.quarantine_dir / (label + "-" + uuid.uuid4().hex)
    os.replace(str(directory), str(destination))


def _cleanup_stale_unverified(paths, snapshot_id: str) -> Optional[str]:
    if not os.path.lexists(str(paths.unverified_path)):
        return None
    try:
        _quarantine_file(
            paths.unverified_path,
            paths.quarantine_dir
            / ("stale-{}-editable-unverified.pptx".format(snapshot_id.split(":", 1)[-1][:12])),
        )
    except (EditableError, OSError) as exc:
        return "stale unverified cleanup failed: {}".format(exc)
    return None


def _promote(
    paths,
    output_lock: OutputLock,
    snapshot_id: str,
    target_kind: str,
    data: bytes,
    result: EditableResult,
    fault_injector,
) -> EditableResult:
    transaction = begin_promotion(
        paths,
        snapshot_id,
        target_kind,
        output_lock=output_lock,
    )
    staged = transaction.directory / "candidate.pptx"
    _write_candidate_bytes(staged, data)
    digest = promote_output(
        transaction,
        staged,
        result,
        fault_injector=fault_injector,
        output_lock=output_lock,
    )
    relative = transaction.target_path.relative_to(paths.run_dir).as_posix()
    promoted = replace(result, output_path=relative, output_sha256=digest)
    if target_kind == "verified":
        cleanup_warning = _cleanup_stale_unverified(paths, snapshot_id)
        if cleanup_warning is not None:
            promoted = replace(
                promoted,
                warnings=promoted.warnings + (cleanup_warning,),
            )
    return promoted


def _default_office_runner(request: Mapping[str, object]) -> OfficeResult:
    script = Path(__file__).resolve().parents[1] / "normalize_and_export.ps1"
    return invoke_office_verification(request, script, 300)


def generate_editable(
    run_dir: Path,
    capability: GenerationCapability,
    fault_injector=None,
) -> EditableResult:
    if not isinstance(capability, GenerationCapability):
        raise TypeError("capability must be GenerationCapability")
    context = None  # type: Optional[RunContext]
    snapshot_id = "sha256:" + "0" * 64
    storyboard = ()
    sources = ()
    try:
        selected = locate_run(Path(run_dir), Path.cwd(), None)
        context = validate_completed_run(selected)
        storyboard = parse_storyboard(context.storyboard_path)
        sources = resolve_slide_sources(context, storyboard)
        config_path = Path(__file__).resolve().parents[2] / "assets" / "verification-config.json"
        config_bytes = config_path.read_bytes()
        config = load_verification_config(config_path)
        snapshot_payload = canonical_snapshot_payload(
            context,
            sources,
            storyboard,
            CONVERTER_VERSION,
            SUBSET_CONTRACT_VERSION,
            config_bytes,
        )
        snapshot_id = compute_snapshot_id(snapshot_payload)
    except EditableError as error:
        return _result(
            "BLOCKED",
            Path(run_dir).name,
            snapshot_id,
            0,
            failures=(_error_failure(error),),
        )
    except (OSError, ValueError) as error:
        return _result(
            "BLOCKED",
            Path(run_dir).name,
            snapshot_id,
            0,
            failures=(_failure("source_unreadable", str(error)),),
        )

    paths = build_output_paths(context.run_dir, context.deck_id)
    try:
        with OutputLock(paths.lock_path) as output_lock:
            recovery = recover_incomplete_transactions(
                paths,
                output_lock=output_lock,
            )
            authoritative_recovery_failures = tuple(
                failure
                for failure in recovery.failures
                if "post-commit cleanup failed" not in failure
            )
            cleanup_recovery_failures = tuple(
                failure
                for failure in recovery.failures
                if "post-commit cleanup failed" in failure
            )
            if authoritative_recovery_failures:
                return _result(
                    "BLOCKED",
                    context.deck_id,
                    snapshot_id,
                    len(storyboard),
                    failures=(
                        _failure(
                            "promotion_conflict",
                            "; ".join(authoritative_recovery_failures),
                        ),
                    ),
                )
            committed = _committed_result(paths, snapshot_id)
            ready_for_verified = capability.office_available and capability.pillow_available
            if committed is not None:
                if committed.status == "PASS":
                    cleanup_warning = _cleanup_stale_unverified(paths, snapshot_id)
                    if cleanup_warning is not None:
                        committed = replace(
                            committed,
                            warnings=committed.warnings + (cleanup_warning,),
                        )
                    return committed
                if committed.status == "GENERATED_UNVERIFIED" and not ready_for_verified:
                    return committed
            if cleanup_recovery_failures:
                return _result(
                    "BLOCKED",
                    context.deck_id,
                    snapshot_id,
                    len(storyboard),
                    failures=(
                        _failure(
                            "promotion_conflict",
                            "; ".join(cleanup_recovery_failures),
                        ),
                    ),
                )

            dependency_failure = _core_dependency_failure()
            if dependency_failure is not None:
                return _result(
                    "BLOCKED",
                    context.deck_id,
                    snapshot_id,
                    len(storyboard),
                    failures=(dependency_failure,),
                )

            try:
                from .svg_parser import DeckPreflightError, preflight_deck

                deck_plan = preflight_deck(
                    context,
                    sources,
                    storyboard,
                    snapshot_id,
                )
            except DeckPreflightError as error:
                return _result(
                    "BLOCKED",
                    context.deck_id,
                    snapshot_id,
                    len(storyboard),
                    failures=tuple(error.failures),
                )
            except EditableError as error:
                return _result(
                    "BLOCKED",
                    context.deck_id,
                    snapshot_id,
                    len(storyboard),
                    failures=(_error_failure(error),),
                )

            work = _create_work_directory(paths)
            keep_work = False
            try:
                try:
                    if committed is not None and committed.status == "GENERATED_UNVERIFIED":
                        candidate_bytes = paths.unverified_path.read_bytes()
                    else:
                        candidate_bytes = _build_presentation_bytes(deck_plan)
                    candidate = work / "candidate.pptx"
                    _write_candidate_bytes(candidate, candidate_bytes)
                except EditableError as error:
                    return _result(
                        "BLOCKED",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=(_error_failure(error),),
                        warnings=tuple(deck_plan.warnings),
                    )
                except (OSError, ValueError) as error:
                    return _result(
                        "BLOCKED",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=(
                            _failure(
                                "candidate_write_failed",
                                str(error),
                            ),
                        ),
                        warnings=tuple(deck_plan.warnings),
                    )
                try:
                    structural = _verify_candidate(candidate, deck_plan, config)
                except Exception:
                    keep_work = True
                    raise
                if not structural.passed:
                    keep_work = True
                    return _result(
                        "FAILED_VERIFICATION",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=tuple(structural.failures),
                        warnings=tuple(deck_plan.warnings),
                    )

                if not ready_for_verified:
                    unverified = _result(
                        "GENERATED_UNVERIFIED",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        warnings=tuple(deck_plan.warnings),
                    )
                    return _promote(
                        paths,
                        output_lock,
                        snapshot_id,
                        "unverified",
                        candidate_bytes,
                        unverified,
                        fault_injector,
                    )

                try:
                    geometry_candidate = work / "geometry-candidate.pptx"
                    _write_candidate_bytes(
                        geometry_candidate,
                        _build_presentation_bytes(deck_plan, include_text=False),
                    )
                    geometry_entries = _write_geometry_svgs(
                        tuple(sources),
                        work / "geometry-svg",
                    )
                except EditableError as error:
                    return _result(
                        "BLOCKED",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=(_error_failure(error),),
                        warnings=tuple(deck_plan.warnings),
                    )
                except (OSError, ValueError) as error:
                    return _result(
                        "BLOCKED",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=(
                            _failure("candidate_write_failed", str(error)),
                        ),
                        warnings=tuple(deck_plan.warnings),
                    )
                try:
                    geometry_failures = _verify_geometry_candidate(
                        geometry_candidate,
                        deck_plan,
                    )
                except Exception:
                    keep_work = True
                    raise
                if geometry_failures:
                    keep_work = True
                    return _result(
                        "FAILED_VERIFICATION",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=tuple(geometry_failures),
                        warnings=tuple(deck_plan.warnings),
                    )
                render_directories = {
                    key: str(work / "renders" / key)
                    for key in RENDER_KEYS
                }
                office_request = {
                    "schema_version": 1,
                    "request_id": uuid.uuid4().hex,
                    "capability_only": False,
                    "protocol_dir": str(work / "office"),
                    "candidate_path": str(candidate),
                    "normalized_path": str(work / "normalized.pptx"),
                    "geometry_candidate_path": str(geometry_candidate),
                    "source_full_deck_path": str(work / "source-full.pptx"),
                    "source_geometry_deck_path": str(work / "source-geometry.pptx"),
                    "selected_svgs": [
                        {"slide_id": source.slide_id, "path": str(source.path)}
                        for source in sources
                    ],
                    "geometry_svgs": list(geometry_entries),
                    "render_directories": render_directories,
                    "ordered_slide_ids": [slide.slide_id for slide in deck_plan.slides],
                    "expected_counts": {
                        "slides": structural.slide_count,
                        "top_level_shapes": structural.top_level_shape_count,
                        "recursive_leaves": structural.recursive_leaf_count,
                        "recursive_groups": structural.recursive_group_count,
                    },
                    "config": {
                        "render_width": config.render_width,
                        "render_height": config.render_height,
                    },
                }
                runner = capability.office_runner or _default_office_runner
                try:
                    office = runner(office_request)
                except Exception:
                    keep_work = True
                    raise
                try:
                    office_binding_failure = _validate_office_result_for_request(
                        office,
                        office_request,
                    )
                except Exception:
                    keep_work = True
                    raise
                if office_binding_failure is not None:
                    keep_work = True
                    return _result(
                        "FAILED_VERIFICATION",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=(office_binding_failure,),
                        warnings=tuple(deck_plan.warnings),
                    )
                if not office.capability:
                    unavailable = bool(
                        office.exit_code == 0
                        and office.error is not None
                        and office.error.get("code") == "powerpoint_unavailable"
                    )
                    if not unavailable:
                        keep_work = True
                        return _result(
                            "FAILED_VERIFICATION",
                            context.deck_id,
                            snapshot_id,
                            len(deck_plan.slides),
                            failures=(_office_failure(office),),
                            warnings=tuple(deck_plan.warnings),
                        )
                    unverified = _result(
                        "GENERATED_UNVERIFIED",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        warnings=tuple(deck_plan.warnings),
                    )
                    return _promote(
                        paths,
                        output_lock,
                        snapshot_id,
                        "unverified",
                        candidate_bytes,
                        unverified,
                        fault_injector,
                    )
                if office.exit_code != 0 or office.error is not None:
                    keep_work = True
                    return _result(
                        "FAILED_VERIFICATION",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=(_office_failure(office),),
                        warnings=tuple(deck_plan.warnings),
                    )
                try:
                    normalized = validate_safe_regular_file(
                        Path(str(office_request["normalized_path"])),
                        (work,),
                    )
                except EditableError as error:
                    keep_work = True
                    return _result(
                        "FAILED_VERIFICATION",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=(
                            _failure(
                                "powerpoint_normalize_failed",
                                error.message,
                            ),
                        ),
                    )
                try:
                    normalized_report = _verify_candidate(
                        normalized,
                        deck_plan,
                        config,
                        office_normalized=True,
                    )
                except Exception:
                    keep_work = True
                    raise
                if not normalized_report.passed:
                    keep_work = True
                    return _result(
                        "FAILED_VERIFICATION",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=tuple(normalized_report.failures),
                        warnings=tuple(deck_plan.warnings),
                    )

                try:
                    from .visual_compare import RenderPaths, compare_render_sets

                    result_render_directories = {
                        key: Path(office.renders[key][0]["path"]).parent
                        for key in RENDER_KEYS
                    }
                    visual = compare_render_sets(
                        RenderPaths(
                            source_full_dir=result_render_directories["source_full"],
                            editable_full_dir=result_render_directories["editable_full"],
                            source_geometry_dir=result_render_directories["source_geometry"],
                            editable_geometry_dir=result_render_directories["editable_geometry"],
                            comparison_dir=work / "comparison",
                        ),
                        tuple(slide.slide_id for slide in deck_plan.slides),
                        config,
                    )
                except Exception:
                    keep_work = True
                    raise
                if not visual.passed:
                    keep_work = True
                    return _result(
                        "FAILED_VERIFICATION",
                        context.deck_id,
                        snapshot_id,
                        len(deck_plan.slides),
                        failures=tuple(visual.failures),
                        warnings=tuple(deck_plan.warnings),
                    )
                passed = _result(
                    "PASS",
                    context.deck_id,
                    snapshot_id,
                    len(deck_plan.slides),
                    warnings=tuple(deck_plan.warnings),
                )
                return _promote(
                    paths,
                    output_lock,
                    snapshot_id,
                    "verified",
                    normalized.read_bytes(),
                    passed,
                    fault_injector,
                )
            finally:
                if work.exists():
                    if keep_work:
                        _quarantine_work(paths, work, "failed-verification")
                    else:
                        shutil.rmtree(work, ignore_errors=True)
    except EditableError as error:
        return _result(
            "BLOCKED",
            context.deck_id,
            snapshot_id,
            len(storyboard),
            failures=(_error_failure(error),),
        )
