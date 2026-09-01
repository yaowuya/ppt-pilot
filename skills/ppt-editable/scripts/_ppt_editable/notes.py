"""Speaker-note formatting and python-pptx note-body persistence."""

from pathlib import Path
from typing import Tuple

from pptx import Presentation

from .contract import StoryboardSlide
from .model import SpeakerNotes


def speaker_notes_from_storyboard(slide: StoryboardSlide) -> SpeakerNotes:
    return SpeakerNotes(
        assertion_title=slide.assertion_title,
        audience_takeaway=slide.audience_takeaway,
        next_link=slide.next_link,
    )


def _note_lines(notes: SpeakerNotes) -> Tuple[str, ...]:
    values = (
        ("本页结论", notes.assertion_title),
        ("听众要点", notes.audience_takeaway),
        ("衔接下一页", notes.next_link),
    )
    return tuple(
        "{}：{}".format(label, value.strip())
        for label, value in values
        if isinstance(value, str) and value.strip()
    )


def format_speaker_notes(notes: SpeakerNotes) -> str:
    return "\n".join(_note_lines(notes))


def attach_speaker_notes(slide, notes: SpeakerNotes) -> None:
    frame = slide.notes_slide.notes_text_frame
    frame.clear()
    lines = _note_lines(notes)
    if not lines:
        return
    frame.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        frame.add_paragraph().text = line


def extract_speaker_notes(presentation_path: Path, slide_index: int) -> Tuple[str, ...]:
    presentation = Presentation(str(presentation_path))
    if type(slide_index) is not int or slide_index < 0 or slide_index >= len(presentation.slides):
        raise IndexError("slide_index out of range")
    frame = presentation.slides[slide_index].notes_slide.notes_text_frame
    return tuple(paragraph.text for paragraph in frame.paragraphs if paragraph.text)
