"""Native DrawingML serialization for preflighted SVG nodes."""

from __future__ import annotations

import hashlib
import io
import json
import math
import re
from typing import Iterable, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement

from .model import Bounds, DeckPlan, SlidePlan, SvgNode, TextLine
from .notes import attach_speaker_notes
from .path_parser import (
    ArcTo,
    ClosePath,
    LineTo,
    MoveTo,
    endpoint_arc_to_center,
    parse_path,
    path_bounds,
    round_int,
)
from .text_layout import choose_primary_font


EMU_PER_PX = 9525
CUSTOM_SCALE = 100
DRAWINGML_COORD_MIN = -27273042329600
DRAWINGML_COORD_MAX = 27273042316900
DRAWINGML_POSITIVE_COORD_MAX = 27273042316900
DRAWINGML_LINE_WIDTH_MAX = 20116800
DRAWINGML_FONT_SIZE_MIN = 100
DRAWINGML_FONT_SIZE_MAX = 400000
DRAWINGML_CHAR_SPACING_MIN = -400000
DRAWINGML_CHAR_SPACING_MAX = 400000
DRAWINGML_ANGLE_MIN = -2147483648
DRAWINGML_ANGLE_MAX = 2147483647
CUSTOM_COORD_MIN = DRAWINGML_COORD_MIN
CUSTOM_COORD_MAX = DRAWINGML_COORD_MAX


class ShapeIdAllocator:
    def __init__(self, first_id: int = 2) -> None:
        if type(first_id) is not int or first_id <= 0:
            raise ValueError("first_id must be a positive integer")
        self._next_id = first_id

    def next(self) -> int:
        value = self._next_id
        if value > 4294967295:
            raise ValueError("shape id exceeds DrawingML range")
        self._next_id += 1
        return value


def trace_description(
    slide_id: str,
    tree_path: str,
    kind: str,
    source_id: Optional[str],
    line_index: Optional[int] = None,
) -> str:
    return json.dumps(
        {
            "schema_version": 1,
            "slide_id": slide_id,
            "tree_path": tree_path,
            "kind": kind,
            "source_id": source_id,
            "line_index": line_index,
        },
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    )


def stable_shape_name(
    slide_id: str,
    tree_path: str,
    kind: str,
    source_id: Optional[str],
    line_index: Optional[int] = None,
) -> str:
    base = source_id or kind
    sanitized = re.sub(r"[^0-9A-Za-z._-]+", "_", base).strip("_") or kind
    sanitized = sanitized[:48]
    identity = "\0".join(
        (slide_id, tree_path, kind, "" if line_index is None else str(line_index))
    ).encode("utf-8")
    suffix = hashlib.sha256(identity).hexdigest()[:8]
    return "{}__{}__{}".format(sanitized, kind, suffix)


def _element(tag_name: str, **attributes):
    node = OxmlElement(tag_name)
    for key, value in attributes.items():
        if value is not None:
            node.set(key, str(value))
    return node


def _checked_int(value: float, minimum: int, maximum: int, label: str) -> int:
    serialized = round_int(value)
    if serialized < minimum or serialized > maximum:
        raise ValueError("{} is outside DrawingML range".format(label))
    return serialized


def _emu(value: float) -> int:
    return _checked_int(
        value * EMU_PER_PX,
        DRAWINGML_COORD_MIN,
        DRAWINGML_COORD_MAX,
        "coordinate",
    )


def _positive_extent(value: float) -> int:
    serialized = _checked_int(
        value * EMU_PER_PX,
        0,
        DRAWINGML_POSITIVE_COORD_MAX,
        "positive extent",
    )
    return max(1, serialized)


def _xfrm(bounds: Bounds, *, group: bool = False, flip_h=False, flip_v=False):
    attributes = {}
    if flip_h:
        attributes["flipH"] = "1"
    if flip_v:
        attributes["flipV"] = "1"
    xfrm = _element("a:xfrm", **attributes)
    off_x = _emu(bounds.left)
    off_y = _emu(bounds.top)
    ext_x = _positive_extent(bounds.right - bounds.left)
    ext_y = _positive_extent(bounds.bottom - bounds.top)
    xfrm.append(_element("a:off", x=off_x, y=off_y))
    xfrm.append(_element("a:ext", cx=ext_x, cy=ext_y))
    if group:
        xfrm.append(_element("a:chOff", x=off_x, y=off_y))
        xfrm.append(_element("a:chExt", cx=ext_x, cy=ext_y))
    return xfrm


def _non_visual_shape(
    shape_id: int,
    name: str,
    description: str,
    *,
    text_box: bool = False,
):
    container = _element("p:nvSpPr")
    container.append(
        _element(
            "p:cNvPr",
            id=shape_id,
            name=name,
            descr=description,
        )
    )
    container.append(_element("p:cNvSpPr", txBox="1" if text_box else None))
    container.append(_element("p:nvPr"))
    return container


def _non_visual_group(shape_id: int, name: str, description: str):
    container = _element("p:nvGrpSpPr")
    container.append(
        _element(
            "p:cNvPr",
            id=shape_id,
            name=name,
            descr=description,
        )
    )
    container.append(_element("p:cNvGrpSpPr"))
    container.append(_element("p:nvPr"))
    return container


def _append_fill(parent, color: Optional[str], opacity: float) -> None:
    if color is None:
        parent.append(_element("a:noFill"))
        return
    fill = _element("a:solidFill")
    color_node = _element("a:srgbClr", val=color.lstrip("#").upper())
    if opacity < 0.0 or opacity > 1.0 or not math.isfinite(opacity):
        raise ValueError("alpha is outside DrawingML range")
    if opacity < 1.0:
        color_node.append(
            _element(
                "a:alpha",
                val=_checked_int(opacity * 100000.0, 0, 100000, "alpha"),
            )
        )
    fill.append(color_node)
    parent.append(fill)


def _append_line(parent, node: SvgNode) -> None:
    line = _element(
        "a:ln",
        w=_checked_int(
            node.style.stroke_width * EMU_PER_PX,
            0,
            DRAWINGML_LINE_WIDTH_MAX,
            "line width",
        ),
    )
    _append_fill(line, node.style.stroke, node.style.stroke_opacity)
    parent.append(line)


def _preset_geometry(name: str):
    geometry = _element("a:prstGeom", prst=name)
    geometry.append(_element("a:avLst"))
    return geometry


def _attributes(node: SvgNode):
    return dict(node.attributes)


def _number(attributes, name, default=0.0):
    value = attributes.get(name)
    return default if value is None else float(value.rstrip("px"))


def _point_list(value: str) -> Tuple[Tuple[float, float], ...]:
    tokens = [token for token in re.split(r"[\s,]+", value.strip()) if token]
    return tuple(
        (float(tokens[index]), float(tokens[index + 1]))
        for index in range(0, len(tokens), 2)
    )


def _primitive_geometry(node: SvgNode):
    attributes = _attributes(node)
    if node.kind == "rect":
        x = _number(attributes, "x")
        y = _number(attributes, "y")
        width = _number(attributes, "width")
        height = _number(attributes, "height")
        return Bounds(x, y, x + width, y + height), _preset_geometry("rect"), False, False
    if node.kind == "circle":
        cx = _number(attributes, "cx")
        cy = _number(attributes, "cy")
        radius = _number(attributes, "r")
        return Bounds(cx - radius, cy - radius, cx + radius, cy + radius), _preset_geometry("ellipse"), False, False
    if node.kind == "ellipse":
        cx = _number(attributes, "cx")
        cy = _number(attributes, "cy")
        rx = _number(attributes, "rx")
        ry = _number(attributes, "ry")
        return Bounds(cx - rx, cy - ry, cx + rx, cy + ry), _preset_geometry("ellipse"), False, False
    if node.kind == "line":
        x1, y1 = _number(attributes, "x1"), _number(attributes, "y1")
        x2, y2 = _number(attributes, "x2"), _number(attributes, "y2")
        return (
            Bounds(min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2)),
            _preset_geometry("line"),
            x2 < x1,
            y2 < y1,
        )
    if node.kind in ("polygon", "polyline"):
        points = _point_list(attributes["points"])
        left = min(point[0] for point in points)
        top = min(point[1] for point in points)
        right = max(point[0] for point in points)
        bottom = max(point[1] for point in points)
        bounds = Bounds(left, top, right, bottom)
        return bounds, _poly_geometry(points, bounds, close=node.kind == "polygon"), False, False
    if node.kind == "path":
        segments = parse_path(attributes["d"])
        bounds = path_bounds(segments)
        return bounds, _path_geometry(segments, bounds), False, False
    raise ValueError("unsupported leaf kind: {}".format(node.kind))


def _custom_container(width: int, height: int, body: Sequence):
    geometry = _element("a:custGeom")
    for name in ("a:avLst", "a:gdLst", "a:ahLst", "a:cxnLst"):
        geometry.append(_element(name))
    geometry.append(_element("a:rect", l=0, t=0, r=width, b=height))
    path_list = _element("a:pathLst")
    path = _element("a:path", w=width, h=height)
    for item in body:
        path.append(item)
    path_list.append(path)
    geometry.append(path_list)
    return geometry


def _relative_point(x: float, y: float, bounds: Bounds):
    return _element(
        "a:pt",
        x=_checked_int(
            (x - bounds.left) * CUSTOM_SCALE,
            CUSTOM_COORD_MIN,
            CUSTOM_COORD_MAX,
            "custom x coordinate",
        ),
        y=_checked_int(
            (y - bounds.top) * CUSTOM_SCALE,
            CUSTOM_COORD_MIN,
            CUSTOM_COORD_MAX,
            "custom y coordinate",
        ),
    )


def _move_or_line(name: str, x: float, y: float, bounds: Bounds):
    command = _element(name)
    command.append(_relative_point(x, y, bounds))
    return command


def _custom_extent(value: float, label: str) -> int:
    return max(
        1,
        _checked_int(
            value * CUSTOM_SCALE,
            0,
            CUSTOM_COORD_MAX,
            label,
        ),
    )


def _poly_geometry(points, bounds: Bounds, *, close: bool):
    body = [_move_or_line("a:moveTo", points[0][0], points[0][1], bounds)]
    body.extend(
        _move_or_line("a:lnTo", x, y, bounds)
        for x, y in points[1:]
    )
    if close:
        body.append(_element("a:close"))
    width = _custom_extent(bounds.right - bounds.left, "custom path width")
    height = _custom_extent(bounds.bottom - bounds.top, "custom path height")
    return _custom_container(width, height, body)


def _path_geometry(segments, bounds: Bounds):
    body = []
    current_x = current_y = None
    subpath_x = subpath_y = None
    for segment in segments:
        if isinstance(segment, MoveTo):
            body.append(_move_or_line("a:moveTo", segment.x, segment.y, bounds))
            current_x, current_y = segment.x, segment.y
            subpath_x, subpath_y = current_x, current_y
        elif isinstance(segment, LineTo):
            body.append(_move_or_line("a:lnTo", segment.x, segment.y, bounds))
            current_x, current_y = segment.x, segment.y
        elif isinstance(segment, ArcTo):
            center = endpoint_arc_to_center(current_x, current_y, segment)
            if center is None:
                body.append(_move_or_line("a:lnTo", segment.end_x, segment.end_y, bounds))
            else:
                start_angle = _checked_int(
                    math.degrees(center.start_radians) * 60000.0,
                    DRAWINGML_ANGLE_MIN,
                    DRAWINGML_ANGLE_MAX,
                    "arc start angle",
                )
                sweep_angle = _checked_int(
                    math.degrees(center.sweep_radians) * 60000.0,
                    DRAWINGML_ANGLE_MIN,
                    DRAWINGML_ANGLE_MAX,
                    "arc sweep angle",
                )
                if center.sweep_radians != 0.0 and sweep_angle == 0:
                    raise ValueError("nonzero SVG arc sweep is below DrawingML precision")
                body.append(
                    _element(
                        "a:arcTo",
                        wR=_checked_int(
                            center.corrected_rx * CUSTOM_SCALE,
                            1,
                            CUSTOM_COORD_MAX,
                            "arc radius x",
                        ),
                        hR=_checked_int(
                            center.corrected_ry * CUSTOM_SCALE,
                            1,
                            CUSTOM_COORD_MAX,
                            "arc radius y",
                        ),
                        stAng=start_angle,
                        swAng=sweep_angle,
                    )
                )
                body.append(_move_or_line("a:lnTo", segment.end_x, segment.end_y, bounds))
            current_x, current_y = segment.end_x, segment.end_y
        else:
            body.append(_element("a:close"))
            current_x, current_y = subpath_x, subpath_y
    width = _custom_extent(bounds.right - bounds.left, "custom path width")
    height = _custom_extent(bounds.bottom - bounds.top, "custom path height")
    return _custom_container(width, height, body)


def build_leaf_shape(node: SvgNode, allocator: ShapeIdAllocator, slide_id: str):
    shape_id = allocator.next()
    source_id = node.style.data_source_id
    name = stable_shape_name(slide_id, node.tree_path, node.kind, source_id)
    description = trace_description(slide_id, node.tree_path, node.kind, source_id)
    bounds, geometry, flip_h, flip_v = _primitive_geometry(node)
    shape = _element("p:sp")
    shape.append(_non_visual_shape(shape_id, name, description))
    properties = _element("p:spPr")
    properties.append(_xfrm(bounds, flip_h=flip_h, flip_v=flip_v))
    properties.append(geometry)
    _append_fill(properties, node.style.fill, node.style.fill_opacity)
    _append_line(properties, node)
    shape.append(properties)
    return shape


def _run_bold(weight: Optional[str]) -> bool:
    if weight == "bold":
        return True
    try:
        return int(weight or "400") >= 600
    except ValueError:
        return False


def build_text_shape(
    line: TextLine,
    allocator: ShapeIdAllocator,
    slide_id: str,
    tree_path: str,
    source_id: Optional[str],
):
    shape_id = allocator.next()
    name = stable_shape_name(
        slide_id,
        tree_path,
        "text",
        source_id,
        line.line_index,
    )
    description = trace_description(
        slide_id,
        tree_path,
        "text",
        source_id,
        line.line_index,
    )
    shape = _element("p:sp")
    shape.append(_non_visual_shape(shape_id, name, description, text_box=True))
    properties = _element("p:spPr")
    properties.append(_xfrm(line.bounds))
    properties.append(_preset_geometry("rect"))
    properties.append(_element("a:noFill"))
    no_line = _element("a:ln")
    no_line.append(_element("a:noFill"))
    properties.append(no_line)
    shape.append(properties)

    text_body = _element("p:txBody")
    body_properties = _element(
        "a:bodyPr",
        wrap="none",
        lIns=0,
        tIns=0,
        rIns=0,
        bIns=0,
        anchor="t",
    )
    body_properties.append(_element("a:noAutofit"))
    text_body.append(body_properties)
    text_body.append(_element("a:lstStyle"))
    paragraph = _element("a:p")
    paragraph.append(
        _element(
            "a:pPr",
            algn={"start": "l", "middle": "ctr", "end": "r"}[line.anchor],
        )
    )
    for run in line.runs:
        font_size = _checked_int(
            (run.style.font_size or 16.0) * 75.0,
            DRAWINGML_FONT_SIZE_MIN,
            DRAWINGML_FONT_SIZE_MAX,
            "text font size",
        )
        character_spacing = (
            _checked_int(
                run.style.letter_spacing * 75.0,
                DRAWINGML_CHAR_SPACING_MIN,
                DRAWINGML_CHAR_SPACING_MAX,
                "text character spacing",
            )
            if run.style.letter_spacing
            else None
        )
        run_node = _element("a:r")
        run_properties = _element(
            "a:rPr",
            lang="zh-CN",
            sz=font_size,
            b="1" if _run_bold(run.style.font_weight) else None,
            spc=character_spacing,
            dirty="0",
        )
        _append_fill(run_properties, run.style.fill or "#000000", run.style.fill_opacity)
        typeface = choose_primary_font(run.style.font_family)
        run_properties.append(_element("a:latin", typeface=typeface))
        run_properties.append(_element("a:ea", typeface=typeface))
        run_properties.append(_element("a:cs", typeface=typeface))
        run_node.append(run_properties)
        text_node = _element("a:t")
        if run.preserve_space or run.text[:1].isspace() or run.text[-1:].isspace():
            text_node.set(
                "{http://www.w3.org/XML/1998/namespace}space",
                "preserve",
            )
        text_node.text = run.text
        run_node.append(text_node)
        paragraph.append(run_node)
    text_body.append(paragraph)
    shape.append(text_body)
    return shape


def _retained_bounds(node: SvgNode, include_text: bool) -> Optional[Bounds]:
    if node.kind == "text" and not include_text:
        return None
    if node.kind != "g":
        return node.bounds
    retained = None
    for child in node.children:
        child_bounds = _retained_bounds(child, include_text)
        if child_bounds is None:
            continue
        retained = child_bounds if retained is None else retained.union(child_bounds)
    if retained is None:
        return None
    if retained.right <= retained.left or retained.bottom <= retained.top:
        raise ValueError("group bounds must have positive extents")
    return retained


def build_group_shape(
    node: SvgNode,
    allocator: ShapeIdAllocator,
    slide_id: str,
    *,
    include_text: bool = True,
):
    bounds = _retained_bounds(node, include_text)
    if bounds is None:
        if include_text:
            raise ValueError("production group has no descendants")
        return None
    shape_id = allocator.next()
    source_id = node.style.data_source_id
    group = _element("p:grpSp")
    group.append(
        _non_visual_group(
            shape_id,
            stable_shape_name(slide_id, node.tree_path, "g", source_id),
            trace_description(slide_id, node.tree_path, "g", source_id),
        )
    )
    properties = _element("p:grpSpPr")
    properties.append(_xfrm(bounds, group=True))
    group.append(properties)
    children = []
    for child in node.children:
        children.extend(
            _build_node(
                child,
                allocator,
                slide_id,
                include_text=include_text,
            )
        )
    if not children:
        raise ValueError("retained group produced no DrawingML children")
    for child in children:
        group.append(child)
    return group


def _build_node(
    node: SvgNode,
    allocator: ShapeIdAllocator,
    slide_id: str,
    *,
    include_text: bool,
):
    if node.kind == "g":
        group = build_group_shape(
            node,
            allocator,
            slide_id,
            include_text=include_text,
        )
        return () if group is None else (group,)
    if node.kind == "text":
        if not include_text:
            return ()
        return tuple(
            build_text_shape(
                line,
                allocator,
                slide_id,
                node.tree_path,
                node.style.data_source_id,
            )
            for line in node.text_lines
        )
    return (build_leaf_shape(node, allocator, slide_id),)


def build_slide(
    slide_plan: SlidePlan,
    allocator: Optional[ShapeIdAllocator] = None,
    *,
    include_text: bool = True,
) -> Tuple:
    allocator = allocator or ShapeIdAllocator()
    result = []
    for node in slide_plan.nodes:
        result.extend(
            _build_node(
                node,
                allocator,
                slide_plan.slide_id,
                include_text=include_text,
            )
        )
    return tuple(result)


def build_presentation(
    deck_plan: DeckPlan,
    *,
    include_text: bool = True,
) -> Presentation:
    if not deck_plan.slides:
        raise ValueError("deck plan has no slides")
    slide_ids = tuple(slide.slide_id for slide in deck_plan.slides)
    if len(slide_ids) != len(set(slide_ids)):
        raise ValueError("duplicate slide owner")
    presentation = Presentation()
    presentation.slide_width = 12192000
    presentation.slide_height = 6858000
    blank_layout = presentation.slide_layouts[6]
    for slide_plan in deck_plan.slides:
        slide = presentation.slides.add_slide(blank_layout)
        slide._element.cSld.set("name", slide_plan.title or slide_plan.slide_id)
        allocator = ShapeIdAllocator()
        for element in build_slide(
            slide_plan,
            allocator,
            include_text=include_text,
        ):
            slide.shapes._spTree.insert_element_before(element, "p:extLst")
        attach_speaker_notes(slide, slide_plan.notes)
    return presentation


def _presentation_stream():
    return io.BytesIO()


def presentation_bytes(
    deck_plan: DeckPlan,
    *,
    include_text: bool = True,
) -> bytes:
    presentation = build_presentation(deck_plan, include_text=include_text)
    with _presentation_stream() as stream:
        presentation.save(stream)
        return stream.getvalue()
