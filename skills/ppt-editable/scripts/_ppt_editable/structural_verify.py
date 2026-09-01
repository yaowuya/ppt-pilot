"""Independent pre-Office verification of editable PPTX candidates."""

from __future__ import annotations

from dataclasses import asdict, dataclass
import io
import json
import math
from pathlib import Path
import posixpath
import re
from typing import Dict, Iterable, Mapping, Optional, Sequence, Tuple
import zipfile

from lxml import etree
from pptx import Presentation

from .drawingml import (
    CUSTOM_COORD_MAX,
    CUSTOM_COORD_MIN,
    DRAWINGML_ANGLE_MAX,
    DRAWINGML_ANGLE_MIN,
    DRAWINGML_COORD_MAX,
    DRAWINGML_COORD_MIN,
    DRAWINGML_POSITIVE_COORD_MAX,
    EMU_PER_PX,
    round_int,
    stable_shape_name,
    trace_description,
)
from .errors import EditableError
from .model import Bounds, DeckPlan, Failure, SvgNode, TextLine, VerificationConfig
from .notes import format_speaker_notes
from .path_parser import (
    ArcTo,
    ClosePath,
    LineTo,
    MoveTo,
    endpoint_arc_to_center,
    parse_path,
    path_bounds,
)
from .text_layout import choose_primary_font


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
RELATIONSHIP_TYPE_PREFIX = "http://schemas.openxmlformats.org/"
OFFICE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
PACKAGE_REL = "http://schemas.openxmlformats.org/package/2006/relationships/"
RELATIONSHIP_TYPES = frozenset(
    {
        OFFICE_REL + name
        for name in (
            "extended-properties",
            "notesMaster",
            "notesSlide",
            "officeDocument",
            "presProps",
            "printerSettings",
            "slide",
            "slideLayout",
            "slideMaster",
            "tableStyles",
            "theme",
            "viewProps",
        )
    }
    | {
        PACKAGE_REL + "metadata/core-properties",
        PACKAGE_REL + "metadata/thumbnail",
    }
)
XML_NS = "http://www.w3.org/XML/1998/namespace"
NS = {"p": P_NS, "a": A_NS, "r": R_NS}
_INTEGER_RE = re.compile(r"^-?[0-9]+$")


@dataclass(frozen=True)
class VerificationReport:
    passed: bool
    failures: Tuple[Failure, ...]
    slide_count: int
    top_level_shape_count: int
    recursive_leaf_count: int
    recursive_group_count: int
    slide_metadata: Tuple[
        Tuple[str, Optional[str], Optional[str]], ...
    ] = ()

    def to_dict(self) -> Mapping[str, object]:
        return {
            "schema_version": 1,
            "kind": "ppt_editable_verification",
            "status": "passed" if self.passed else "failed",
            "slide_count": self.slide_count,
            "top_level_shape_count": self.top_level_shape_count,
            "recursive_leaf_count": self.recursive_leaf_count,
            "recursive_group_count": self.recursive_group_count,
            "slides": [
                {
                    "slide_id": slide_id,
                    "title": title,
                    "description": description,
                }
                for slide_id, title, description in self.slide_metadata
            ],
            "failures": [asdict(failure) for failure in self.failures],
        }


def _failure(
    code: str,
    message: str,
    *,
    slide_id: Optional[str] = None,
    tree_path: Optional[str] = None,
    element_type: Optional[str] = None,
) -> Failure:
    return Failure(
        code=code,
        slide_id=slide_id,
        svg_tree_path=tree_path,
        element_type=element_type,
        message=message,
        remediation="rebuild the candidate from the unchanged authoritative SVG plan",
    )


def _xml(data: bytes):
    parser = etree.XMLParser(resolve_entities=False, no_network=True, recover=False)
    return etree.fromstring(data, parser=parser)


def _zip_payload(path: Path):
    try:
        data = Path(path).read_bytes()
        archive = zipfile.ZipFile(io.BytesIO(data), "r")
    except (OSError, zipfile.BadZipFile) as exc:
        return None, (
            _failure("pptx_zip_invalid", "candidate is not a readable ZIP package"),
        )
    with archive:
        infos = archive.infolist()
        names = [info.filename for info in infos]
        failures = []
        if len(names) != len(set(names)):
            failures.append(_failure("pptx_zip_invalid", "candidate ZIP has duplicate entries"))
        try:
            bad_entry = archive.testzip()
        except Exception:
            failures.append(_failure("pptx_zip_invalid", "candidate ZIP integrity scan failed"))
            return {}, tuple(failures)
        if bad_entry is not None:
            failures.append(_failure("pptx_zip_invalid", "candidate ZIP CRC check failed"))
        required = {
            "[Content_Types].xml",
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
        }
        if not required.issubset(names):
            failures.append(_failure("pptx_zip_invalid", "candidate ZIP lacks required package parts"))
        payload = {}
        for info in infos:
            if info.filename not in payload:
                try:
                    payload[info.filename] = archive.read(info)
                except Exception:
                    failures.append(_failure("pptx_zip_invalid", "candidate ZIP entry cannot be read"))
        return payload, tuple(failures)


def _relationship_target_part(rels_name: str, target: str) -> Optional[str]:
    if target.startswith("/"):
        normalized = posixpath.normpath(target.lstrip("/"))
    else:
        marker = "/_rels/"
        if marker in rels_name:
            prefix, rel_name = rels_name.split(marker, 1)
            source_name = rel_name[: -len(".rels")]
            source_parent = posixpath.dirname(posixpath.join(prefix, source_name))
        elif rels_name == "_rels/.rels":
            source_parent = ""
        else:
            return None
        normalized = posixpath.normpath(posixpath.join(source_parent, target))
    if normalized == ".." or normalized.startswith("../") or normalized.startswith("/"):
        return None
    return normalized


def _verify_content_types(payload: Mapping[str, bytes]) -> Tuple[Failure, ...]:
    try:
        root = _xml(payload["[Content_Types].xml"])
    except (KeyError, etree.XMLSyntaxError):
        return (_failure("structure_mismatch", "content-types part is malformed"),)
    if root.tag != "{{{}}}Types".format(CT_NS):
        return (_failure("structure_mismatch", "content-types namespace is invalid"),)
    failures = []
    if root.attrib or (root.text or "").strip():
        failures.append(
            _failure(
                "structure_mismatch",
                "content-types root structure is invalid",
            )
        )
    defaults = {}
    overrides = {}
    for child in root:
        if child.tag == "{{{}}}Default".format(CT_NS):
            extension = child.get("Extension")
            content_type = child.get("ContentType")
            extension_key = extension.lower() if isinstance(extension, str) else None
            if (
                set(child.attrib) != {"Extension", "ContentType"}
                or len(child)
                or (child.text or "").strip()
                or (child.tail or "").strip()
                or not extension_key
                or not content_type
                or extension_key in defaults
            ):
                failures.append(_failure("structure_mismatch", "default content-type declaration is invalid"))
            else:
                defaults[extension_key] = content_type
        elif child.tag == "{{{}}}Override".format(CT_NS):
            part_name = child.get("PartName")
            content_type = child.get("ContentType")
            normalized_part = part_name.lstrip("/") if isinstance(part_name, str) else None
            if (
                set(child.attrib) != {"PartName", "ContentType"}
                or len(child)
                or (child.text or "").strip()
                or (child.tail or "").strip()
                or not part_name
                or not part_name.startswith("/")
                or not content_type
                or normalized_part in overrides
                or normalized_part == "[Content_Types].xml"
                or posixpath.normpath(normalized_part or "") != normalized_part
                or normalized_part not in payload
            ):
                failures.append(_failure("structure_mismatch", "override content-type declaration is invalid"))
            else:
                overrides[normalized_part] = content_type
        else:
            failures.append(_failure("structure_mismatch", "unexpected content-types element"))
    if defaults.get("rels") != "application/vnd.openxmlformats-package.relationships+xml":
        failures.append(_failure("structure_mismatch", "canonical rels content type default is missing"))
    for name in payload:
        if name == "[Content_Types].xml":
            continue
        if name.endswith(".rels"):
            expected_rels_type = (
                "application/vnd.openxmlformats-package.relationships+xml"
            )
            effective = overrides.get(name, defaults.get("rels"))
            if effective != expected_rels_type:
                failures.append(
                    _failure(
                        "structure_mismatch",
                        "relationship part lacks canonical content type",
                    )
                )
            continue
        extension = name.rsplit(".", 1)[-1].lower() if "." in name else ""
        if name not in overrides and extension not in defaults:
            failures.append(_failure("structure_mismatch", "package part lacks a content type: {}".format(name)))
    expected = {
        "ppt/presentation.xml": "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
        "ppt/presProps.xml": "application/vnd.openxmlformats-officedocument.presentationml.presProps+xml",
        "ppt/viewProps.xml": "application/vnd.openxmlformats-officedocument.presentationml.viewProps+xml",
        "ppt/tableStyles.xml": "application/vnd.openxmlformats-officedocument.presentationml.tableStyles+xml",
        "docProps/core.xml": "application/vnd.openxmlformats-package.core-properties+xml",
        "docProps/app.xml": "application/vnd.openxmlformats-officedocument.extended-properties+xml",
    }
    patterns = (
        (r"ppt/slides/slide[0-9]+\.xml", "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"),
        (r"ppt/notesSlides/notesSlide[0-9]+\.xml", "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"),
        (r"ppt/slideLayouts/slideLayout[0-9]+\.xml", "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"),
        (r"ppt/slideMasters/slideMaster[0-9]+\.xml", "application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"),
        (r"ppt/notesMasters/notesMaster[0-9]+\.xml", "application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"),
        (r"ppt/theme/theme[0-9]+\.xml", "application/vnd.openxmlformats-officedocument.theme+xml"),
        (r"ppt/printerSettings/printerSettings[0-9]+\.bin", "application/vnd.openxmlformats-officedocument.presentationml.printerSettings"),
        (r"docProps/thumbnail\.(?:jpeg|jpg)", "image/jpeg"),
    )
    for name in payload:
        for pattern, content_type in patterns:
            if re.fullmatch(pattern, name):
                expected[name] = content_type
                break
    for name, content_type in expected.items():
        if name not in payload:
            continue
        extension = name.rsplit(".", 1)[-1].lower() if "." in name else ""
        effective = overrides.get(name, defaults.get(extension))
        if effective != content_type:
            failures.append(_failure("structure_mismatch", "content type differs for {}".format(name)))
    return tuple(failures)


def _ordered_slide_parts(
    payload: Mapping[str, bytes],
) -> Tuple[Tuple[str, ...], Tuple[Failure, ...]]:
    try:
        presentation = _xml(payload["ppt/presentation.xml"])
        relationships = _xml(payload["ppt/_rels/presentation.xml.rels"])
    except (KeyError, etree.XMLSyntaxError):
        return (), (_failure("structure_mismatch", "presentation order parts are malformed"),)
    if presentation.tag != "{{{}}}presentation".format(P_NS):
        return (), (_failure("structure_mismatch", "presentation root namespace is invalid"),)
    if relationships.tag != "{{{}}}Relationships".format(R_NS):
        return (), (_failure("structure_mismatch", "presentation relationships root is invalid"),)
    relationship_map = {}
    for relationship in relationships:
        if relationship.tag != "{{{}}}Relationship".format(R_NS):
            return (), (_failure("structure_mismatch", "unexpected presentation relationship element"),)
        relationship_id = relationship.get("Id")
        if not relationship_id or relationship_id in relationship_map:
            return (), (_failure("structure_mismatch", "presentation relationship IDs are invalid"),)
        relationship_map[relationship_id] = relationship
    slide_parts = []
    seen_relationships = set()
    for slide_id in presentation.findall("./p:sldIdLst/p:sldId", NS):
        relationship_id = slide_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        relationship = relationship_map.get(relationship_id)
        if relationship is None or relationship_id in seen_relationships:
            return (), (_failure("structure_mismatch", "slide relationship order is invalid"),)
        seen_relationships.add(relationship_id)
        if relationship.get("Type") != "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide":
            return (), (_failure("structure_mismatch", "presentation slide relationship type is invalid"),)
        target = relationship.get("Target")
        part = _relationship_target_part(
            "ppt/_rels/presentation.xml.rels",
            target,
        ) if isinstance(target, str) else None
        if part is None or part not in payload or part in slide_parts:
            return (), (_failure("structure_mismatch", "presentation slide target is invalid"),)
        slide_parts.append(part)
    if not slide_parts:
        return (), (_failure("structure_mismatch", "presentation has no ordered slides"),)
    return tuple(slide_parts), ()


def _relationship_source_part(rels_name: str) -> Optional[str]:
    if rels_name == "_rels/.rels":
        return ""
    marker = "/_rels/"
    if marker not in rels_name or not rels_name.endswith(".rels"):
        return None
    prefix, relationship_name = rels_name.split(marker, 1)
    return posixpath.join(prefix, relationship_name[: -len(".rels")])


def _part_class(name: str) -> Optional[str]:
    patterns = (
        (r"", "root"),
        (r"ppt/presentation\.xml", "presentation"),
        (r"ppt/presProps\.xml", "presProps"),
        (r"ppt/viewProps\.xml", "viewProps"),
        (r"ppt/tableStyles\.xml", "tableStyles"),
        (r"ppt/slides/slide[0-9]+\.xml", "slide"),
        (r"ppt/slideLayouts/slideLayout[0-9]+\.xml", "slideLayout"),
        (r"ppt/slideMasters/slideMaster[0-9]+\.xml", "slideMaster"),
        (r"ppt/notesSlides/notesSlide[0-9]+\.xml", "notesSlide"),
        (r"ppt/notesMasters/notesMaster[0-9]+\.xml", "notesMaster"),
        (r"ppt/theme/theme[0-9]+\.xml", "theme"),
        (r"ppt/printerSettings/printerSettings[0-9]+\.bin", "printerSettings"),
        (r"docProps/core\.xml", "coreProps"),
        (r"docProps/app\.xml", "appProps"),
        (r"docProps/thumbnail\.(?:jpeg|jpg)", "thumbnail"),
    )
    return next((kind for pattern, kind in patterns if re.fullmatch(pattern, name)), None)


_ALLOWED_RELATIONSHIP_TRIPLES = frozenset(
    {
        ("root", OFFICE_REL + "officeDocument", "presentation"),
        ("root", PACKAGE_REL + "metadata/core-properties", "coreProps"),
        ("root", OFFICE_REL + "extended-properties", "appProps"),
        ("root", PACKAGE_REL + "metadata/thumbnail", "thumbnail"),
        ("presentation", OFFICE_REL + "slideMaster", "slideMaster"),
        ("presentation", OFFICE_REL + "printerSettings", "printerSettings"),
        ("presentation", OFFICE_REL + "presProps", "presProps"),
        ("presentation", OFFICE_REL + "viewProps", "viewProps"),
        ("presentation", OFFICE_REL + "theme", "theme"),
        ("presentation", OFFICE_REL + "tableStyles", "tableStyles"),
        ("presentation", OFFICE_REL + "slide", "slide"),
        ("presentation", OFFICE_REL + "notesMaster", "notesMaster"),
        ("slideMaster", OFFICE_REL + "slideLayout", "slideLayout"),
        ("slideMaster", OFFICE_REL + "theme", "theme"),
        ("slideLayout", OFFICE_REL + "slideMaster", "slideMaster"),
        ("slide", OFFICE_REL + "slideLayout", "slideLayout"),
        ("slide", OFFICE_REL + "notesSlide", "notesSlide"),
        ("notesSlide", OFFICE_REL + "notesMaster", "notesMaster"),
        ("notesSlide", OFFICE_REL + "slide", "slide"),
        ("notesMaster", OFFICE_REL + "theme", "theme"),
    }
)


_RELATIONSHIP_OWNER_CLASSES = frozenset(
    {
        "root",
        "presentation",
        "slideMaster",
        "slideLayout",
        "slide",
        "notesSlide",
        "notesMaster",
    }
)


def _verify_relationship_graph(
    payload: Mapping[str, bytes],
    *,
    office_normalized: bool = False,
) -> Tuple[Failure, ...]:
    failures = []
    edges = {}
    for name, data in payload.items():
        if not name.endswith(".rels"):
            continue
        source = _relationship_source_part(name)
        if source is None or (source and source not in payload):
            failures.append(_failure("structure_mismatch", "relationship part has no source: {}".format(name)))
            continue
        source_class = _part_class(source)
        if source_class not in _RELATIONSHIP_OWNER_CLASSES:
            failures.append(
                _failure(
                    "structure_mismatch",
                    "relationship source cannot own relationships: {}".format(name),
                )
            )
            continue
        try:
            root = _xml(data)
        except etree.XMLSyntaxError:
            failures.append(_failure("structure_mismatch", "relationship part is malformed: {}".format(name)))
            continue
        if root.tag != "{{{}}}Relationships".format(R_NS):
            failures.append(_failure("structure_mismatch", "relationship namespace is invalid: {}".format(name)))
            continue
        if root.attrib or (root.text or "").strip():
            failures.append(
                _failure(
                    "structure_mismatch",
                    "relationship root structure is invalid: {}".format(name),
                )
            )
        seen_ids = set()
        source_edges = []
        for relationship in root:
            if relationship.tag != "{{{}}}Relationship".format(R_NS):
                failures.append(_failure("structure_mismatch", "unexpected relationship element: {}".format(name)))
                continue
            required = {"Id", "Type", "Target"}
            allowed = required | {"TargetMode"}
            relationship_id = relationship.get("Id")
            relationship_type = relationship.get("Type")
            target = relationship.get("Target")
            if (
                not required.issubset(relationship.attrib)
                or set(relationship.attrib) - allowed
                or len(relationship)
                or (relationship.text or "").strip()
                or (relationship.tail or "").strip()
                or not relationship_id
                or relationship_id in seen_ids
                or not isinstance(relationship_type, str)
                or relationship_type not in RELATIONSHIP_TYPES
                or not isinstance(target, str)
            ):
                failures.append(_failure("structure_mismatch", "relationship declaration is invalid: {}".format(name)))
                continue
            seen_ids.add(relationship_id)
            target_mode = relationship.get("TargetMode")
            if target_mode not in (None, "Internal"):
                if target_mode == "External":
                    failures.append(
                        _failure(
                            "image_fallback_detected",
                            "external relationship is forbidden",
                        )
                    )
                else:
                    failures.append(
                        _failure(
                            "structure_mismatch",
                            "relationship TargetMode is invalid: {}".format(name),
                        )
                    )
                continue
            resolved = _relationship_target_part(name, target)
            if resolved is None or resolved not in payload:
                failures.append(_failure("structure_mismatch", "relationship target is missing or unsafe in {}".format(name)))
                continue
            target_class = _part_class(resolved)
            if (
                target_class is None
                or (source_class, relationship_type, target_class)
                not in _ALLOWED_RELATIONSHIP_TRIPLES
            ):
                failures.append(
                    _failure(
                        "structure_mismatch",
                        "relationship source/type/target role is invalid in {}".format(name),
                    )
                )
                continue
            source_edges.append((relationship_type, resolved, relationship_id))
        edges[source] = tuple(source_edges)

    def targets_for_type(source, relationship_type):
        return tuple(
            target
            for rel_type, target, _ in edges.get(source, ())
            if rel_type == relationship_type
        )

    def targets(source, role):
        return targets_for_type(source, OFFICE_REL + role.lstrip("/"))

    root_requirements = (
        (OFFICE_REL + "officeDocument", "ppt/presentation.xml"),
        (PACKAGE_REL + "metadata/core-properties", "docProps/core.xml"),
        (OFFICE_REL + "extended-properties", "docProps/app.xml"),
        (PACKAGE_REL + "metadata/thumbnail", "docProps/thumbnail.jpeg"),
    )
    for relationship_type, expected_target in root_requirements:
        if targets_for_type("", relationship_type) != (expected_target,):
            failures.append(
                _failure(
                    "structure_mismatch",
                    "package root relationship cardinality differs for {}".format(
                        relationship_type.rsplit("/", 1)[-1]
                    ),
                )
            )
    presentation_requirements = (
        (OFFICE_REL + "printerSettings", "ppt/printerSettings/printerSettings1.bin"),
        (OFFICE_REL + "presProps", "ppt/presProps.xml"),
        (OFFICE_REL + "viewProps", "ppt/viewProps.xml"),
        (OFFICE_REL + "theme", "ppt/theme/theme1.xml"),
        (OFFICE_REL + "tableStyles", "ppt/tableStyles.xml"),
    )
    for relationship_type, expected_target in presentation_requirements:
        observed_targets = targets_for_type(
            "ppt/presentation.xml",
            relationship_type,
        )
        may_be_removed_by_office = bool(
            office_normalized
            and relationship_type == OFFICE_REL + "printerSettings"
            and not observed_targets
        )
        if observed_targets != (expected_target,) and not may_be_removed_by_office:
            failures.append(
                _failure(
                    "structure_mismatch",
                    "presentation support relationship cardinality differs for {}".format(
                        relationship_type.rsplit("/", 1)[-1]
                    ),
                )
            )
    presentation_slides = targets("ppt/presentation.xml", "slide")
    slide_parts = {
        name
        for name in payload
        if re.fullmatch(r"ppt/slides/slide[0-9]+\.xml", name)
    }
    if set(presentation_slides) != slide_parts or len(presentation_slides) != len(slide_parts):
        failures.append(_failure("structure_mismatch", "slide part inventory differs from presentation order"))
    expected_notes = set()
    for slide in presentation_slides:
        layouts = targets(slide, "/slideLayout")
        notes = targets(slide, "/notesSlide")
        expected_notes.update(notes)
        if len(layouts) != 1 or len(notes) != 1:
            failures.append(_failure("structure_mismatch", "slide relationships are incomplete: {}".format(slide)))
        for note in notes:
            if targets(note, "/slide") != (slide,) or len(targets(note, "/notesMaster")) != 1:
                failures.append(_failure("structure_mismatch", "notes relationship reciprocity differs: {}".format(note)))
    note_parts = {
        name
        for name in payload
        if re.fullmatch(r"ppt/notesSlides/notesSlide[0-9]+\.xml", name)
    }
    if expected_notes != note_parts:
        failures.append(_failure("structure_mismatch", "notes part inventory differs from slide relationships"))
    slide_master_parts = {
        name
        for name in payload
        if re.fullmatch(r"ppt/slideMasters/slideMaster[0-9]+\.xml", name)
    }
    presentation_masters = targets("ppt/presentation.xml", "/slideMaster")
    if (
        not slide_master_parts
        or set(presentation_masters) != slide_master_parts
        or len(presentation_masters) != len(slide_master_parts)
    ):
        failures.append(
            _failure(
                "structure_mismatch",
                "slide-master inventory differs from presentation relationships",
            )
        )

    layout_parts = {
        name
        for name in payload
        if re.fullmatch(r"ppt/slideLayouts/slideLayout[0-9]+\.xml", name)
    }
    master_layouts = []
    for master in slide_master_parts:
        layouts = targets(master, "/slideLayout")
        master_layouts.extend(layouts)
        if not layouts or len(targets(master, "/theme")) != 1:
            failures.append(
                _failure(
                    "structure_mismatch",
                    "slide master relationships are incomplete: {}".format(master),
                )
            )
    if set(master_layouts) != layout_parts or len(master_layouts) != len(layout_parts):
        failures.append(
            _failure(
                "structure_mismatch",
                "slide-layout inventory differs from master relationships",
            )
        )
    for layout in layout_parts:
        if len(targets(layout, "/slideMaster")) != 1:
            failures.append(
                _failure(
                    "structure_mismatch",
                    "slide layout lacks one master relationship",
                )
            )

    notes_master_parts = {
        name
        for name in payload
        if re.fullmatch(r"ppt/notesMasters/notesMaster[0-9]+\.xml", name)
    }
    presentation_notes_masters = targets("ppt/presentation.xml", "/notesMaster")
    if (
        set(presentation_notes_masters) != notes_master_parts
        or len(presentation_notes_masters) != len(notes_master_parts)
    ):
        failures.append(
            _failure(
                "structure_mismatch",
                "notes-master inventory differs from presentation relationships",
            )
        )
    for notes_master in notes_master_parts:
        if len(targets(notes_master, "/theme")) != 1:
            failures.append(
                _failure(
                    "structure_mismatch",
                    "notes master lacks one theme relationship",
                )
            )

    reachable = set()
    pending = [""]
    while pending:
        source = pending.pop()
        if source in reachable:
            continue
        reachable.add(source)
        for _, target, _ in edges.get(source, ()):
            if target not in reachable:
                pending.append(target)
    for name in payload:
        if name == "[Content_Types].xml" or name.endswith(".rels"):
            continue
        if name not in reachable:
            failures.append(_failure("structure_mismatch", "package part is unreachable: {}".format(name)))
    return tuple(failures)


_ALLOWED_NOTES_PLACEHOLDER_TYPES = frozenset(
    {"sldImg", "body", "dt", "hdr", "ftr", "sldNum"}
)
_REQUIRED_NOTES_PLACEHOLDER_TYPES = frozenset({"sldImg", "body"})


def _verify_notes_shape_inventory(
    payload: Mapping[str, bytes],
) -> Tuple[Failure, ...]:
    failures = []
    notes_parts = sorted(
        name
        for name in payload
        if re.fullmatch(r"ppt/notesSlides/notesSlide[0-9]+\.xml", name)
    )
    for name in notes_parts:
        try:
            root = _xml(payload[name])
        except etree.XMLSyntaxError:
            failures.append(
                _failure("structure_mismatch", "notes slide is malformed: {}".format(name))
            )
            continue
        if root.tag != "{{{}}}notes".format(P_NS):
            failures.append(
                _failure("structure_mismatch", "notes slide root differs: {}".format(name))
            )
            continue
        trees = root.findall("./p:cSld/p:spTree", NS)
        if len(trees) != 1:
            failures.append(
                _failure("structure_mismatch", "notes slide shape tree differs: {}".format(name))
            )
            continue
        tree = trees[0]
        children = list(tree)
        if len(children) < 2 or [child.tag for child in children[:2]] != [
            "{{{}}}nvGrpSpPr".format(P_NS),
            "{{{}}}grpSpPr".format(P_NS),
        ]:
            failures.append(
                _failure("structure_mismatch", "notes slide shape tree root differs: {}".format(name))
            )
            continue
        shape_children = children[2:]
        if shape_children and shape_children[-1].tag == "{{{}}}extLst".format(P_NS):
            shape_children = shape_children[:-1]
        counts = {}
        inventory_valid = True
        for shape in shape_children:
            if shape.tag != "{{{}}}sp".format(P_NS):
                inventory_valid = False
                continue
            placeholders = shape.findall("./p:nvSpPr/p:nvPr/p:ph", NS)
            if len(placeholders) != 1:
                inventory_valid = False
                continue
            placeholder_type = placeholders[0].get("type")
            if placeholder_type not in _ALLOWED_NOTES_PLACEHOLDER_TYPES:
                inventory_valid = False
                continue
            counts[placeholder_type] = counts.get(placeholder_type, 0) + 1
        if (
            not inventory_valid
            or any(counts.get(kind) != 1 for kind in _REQUIRED_NOTES_PLACEHOLDER_TYPES)
            or any(count > 1 for count in counts.values())
        ):
            failures.append(
                _failure(
                    "structure_mismatch",
                    "notes slide placeholder inventory differs: {}".format(name),
                )
            )
    return tuple(failures)


def _verify_payload_parts(
    payload: Mapping[str, bytes],
    initial_failures: Sequence[Failure] = (),
    *,
    office_normalized: bool = False,
) -> Tuple[Failure, ...]:
    result = list(initial_failures)
    if any(name.startswith("ppt/media/") for name in payload):
        result.append(_failure("image_fallback_detected", "candidate contains media parts"))
    result.extend(
        _verify_relationship_graph(
            payload,
            office_normalized=office_normalized,
        )
    )
    result.extend(_verify_content_types(payload))
    result.extend(_verify_notes_shape_inventory(payload))
    return tuple(result)


def verify_zip_and_parts(pptx_path: Path) -> Tuple[Failure, ...]:
    payload, failures = _zip_payload(pptx_path)
    if payload is None:
        return failures
    return _verify_payload_parts(payload, failures)


def _expected_trace(slide_id: str, node: SvgNode, line: Optional[TextLine] = None):
    return json.loads(
        trace_description(
            slide_id,
            node.tree_path,
            "text" if line is not None else node.kind,
            node.style.data_source_id,
            None if line is None else line.line_index,
        )
    )


def _expected_tree(
    slide_id: str,
    nodes: Sequence[SvgNode],
    next_id=None,
):
    if next_id is None:
        next_id = [2]
    result = []
    for node in nodes:
        if node.kind == "g":
            shape_id = next_id[0]
            next_id[0] += 1
            result.append(
                {
                    "id": shape_id,
                    "name": stable_shape_name(
                        slide_id,
                        node.tree_path,
                        "g",
                        node.style.data_source_id,
                    ),
                    "trace": _expected_trace(slide_id, node),
                    "children": _expected_tree(slide_id, node.children, next_id),
                }
            )
        elif node.kind == "text":
            for line in node.text_lines:
                shape_id = next_id[0]
                next_id[0] += 1
                result.append(
                    {
                        "id": shape_id,
                        "name": stable_shape_name(
                            slide_id,
                            node.tree_path,
                            "text",
                            node.style.data_source_id,
                            line.line_index,
                        ),
                        "trace": _expected_trace(slide_id, node, line),
                        "line": line,
                    }
                )
        else:
            shape_id = next_id[0]
            next_id[0] += 1
            result.append(
                {
                    "id": shape_id,
                    "name": stable_shape_name(
                        slide_id,
                        node.tree_path,
                        node.kind,
                        node.style.data_source_id,
                    ),
                    "trace": _expected_trace(slide_id, node),
                }
            )
    return result


def _shape_c_nv_pr(element):
    if element.tag == "{{{}}}sp".format(P_NS):
        return element.find("./p:nvSpPr/p:cNvPr", NS)
    if element.tag == "{{{}}}grpSp".format(P_NS):
        return element.find("./p:nvGrpSpPr/p:cNvPr", NS)
    return None


def _trace(element):
    node = _shape_c_nv_pr(element)
    if node is None or node.get("descr") is None:
        return None
    try:
        value = json.loads(node.get("descr"))
    except (json.JSONDecodeError, ValueError, RecursionError, MemoryError):
        return None
    return value if isinstance(value, dict) else None


def _shape_identity(element):
    node = _shape_c_nv_pr(element)
    if node is None:
        return None, None, None
    raw_id = node.get("id")
    if (
        raw_id is None
        or len(raw_id) > 10
        or re.fullmatch(r"[0-9]+", raw_id) is None
    ):
        shape_id = None
    else:
        shape_id = int(raw_id)
        if shape_id < 0 or shape_id > 4294967295:
            shape_id = None
    return shape_id, node.get("name"), _trace(element)


def _shape_children(container):
    children = list(container[2:])
    if children and children[-1].tag == "{{{}}}extLst".format(P_NS):
        children.pop()
    return tuple(children)


def _actual_tree(elements, counters):
    result = []
    for element in elements:
        if element.tag == "{{{}}}pic".format(P_NS):
            counters["pics"] += 1
            result.append({"picture": True})
            continue
        shape_id, name, trace = _shape_identity(element)
        if element.tag == "{{{}}}grpSp".format(P_NS):
            counters["groups"] += 1
            result.append(
                {
                    "id": shape_id,
                    "name": name,
                    "trace": trace,
                    "children": _actual_tree(_shape_children(element), counters),
                }
            )
        else:
            counters["leaves"] += 1
            result.append(
                {
                    "id": shape_id,
                    "name": name,
                    "trace": trace,
                    "element": element,
                }
            )
    return result


def _identity_only(value):
    result = []
    for item in value:
        converted = {
            "id": item.get("id"),
            "name": item.get("name"),
            "trace": item.get("trace"),
        }
        if "children" in item:
            converted["children"] = _identity_only(item["children"])
        result.append(converted)
    return result


def _expected_text_lines(nodes: Sequence[SvgNode]) -> Dict[Tuple[str, int], TextLine]:
    result = {}
    for node in nodes:
        if node.kind == "g":
            result.update(_expected_text_lines(node.children))
        elif node.kind == "text":
            for line in node.text_lines:
                result[(node.tree_path, line.line_index)] = line
    return result


def _exact_solid_fill_matches(parent, expected_color: str, expected_alpha: int) -> bool:
    fills = parent.findall("./a:solidFill", NS)
    if len(fills) != 1 or fills[0].attrib or len(fills[0]) != 1:
        return False
    color = fills[0][0]
    if (
        color.tag != "{{{}}}srgbClr".format(A_NS)
        or set(color.attrib) != {"val"}
        or color.get("val") != expected_color
    ):
        return False
    if expected_alpha == 100000:
        return len(color) == 0
    if len(color) != 1:
        return False
    alpha = color[0]
    return bool(
        alpha.tag == "{{{}}}alpha".format(A_NS)
        and set(alpha.attrib) == {"val"}
        and not len(alpha)
        and alpha.get("val") == str(expected_alpha)
    )


def _verify_text_shape(
    element,
    line: TextLine,
    slide_id: str,
    *,
    office_normalized: bool = False,
) -> Tuple[Failure, ...]:
    failures = []
    shape_properties = element.find("./p:spPr", NS)
    line_properties = (
        None
        if shape_properties is None
        else shape_properties.find("./a:ln", NS)
    )
    if (
        shape_properties is None
        or not _exact_no_fill_matches(shape_properties)
        or line_properties is None
        or line_properties.attrib
        or [child.tag for child in line_properties]
        != ["{{{}}}noFill".format(A_NS)]
        or not _exact_no_fill_matches(line_properties)
    ):
        failures.append(
            _failure(
                "content_mismatch",
                "text box fill or border differs",
                slide_id=slide_id,
            )
        )
    c_nv_sp = element.find("./p:nvSpPr/p:cNvSpPr", NS)
    body = element.find("./p:txBody", NS)
    if c_nv_sp is None or c_nv_sp.get("txBox") != "1" or body is None:
        return (_failure("content_mismatch", "text trace is not a true text box", slide_id=slide_id),)
    body_properties = body.find("./a:bodyPr", NS)
    list_style = body.find("./a:lstStyle", NS)
    if (
        body_properties is None
        or set(body_properties.attrib)
        != {"wrap", "lIns", "tIns", "rIns", "bIns", "anchor"}
        or body_properties.get("wrap") != "none"
        or body_properties.get("anchor") != "t"
        or any(body_properties.get(name) != "0" for name in ("lIns", "tIns", "rIns", "bIns"))
        or [child.tag for child in body_properties] != ["{{{}}}noAutofit".format(A_NS)]
        or list_style is None
        or list_style.attrib
        or len(list_style)
    ):
        failures.append(
            _failure(
                "content_mismatch",
                "text body properties differ",
                slide_id=slide_id,
            )
        )
    expected_body_children = [
        "{{{}}}bodyPr".format(A_NS),
        "{{{}}}lstStyle".format(A_NS),
        "{{{}}}p".format(A_NS),
    ]
    if [child.tag for child in body] != expected_body_children:
        failures.append(_failure("content_mismatch", "text body child structure differs", slide_id=slide_id))
    paragraphs = body.findall("./a:p", NS)
    if len(paragraphs) != 1:
        failures.append(_failure("content_mismatch", "text box must contain exactly one paragraph", slide_id=slide_id))
        return tuple(failures)
    paragraph = paragraphs[0]
    if (
        not len(paragraph)
        or paragraph[0].tag != "{{{}}}pPr".format(A_NS)
        or any(child.tag != "{{{}}}r".format(A_NS) for child in paragraph[1:])
    ):
        failures.append(_failure("content_mismatch", "text paragraph structure differs", slide_id=slide_id))
        return tuple(failures)
    actual_runs = paragraph.findall("./a:r", NS)
    if len(actual_runs) != len(line.runs):
        failures.append(_failure("content_mismatch", "text run count differs", slide_id=slide_id))
        return tuple(failures)
    alignment = paragraph.find("./a:pPr", NS)
    expected_alignment = {"start": "l", "middle": "ctr", "end": "r"}[line.anchor]
    if (
        alignment is None
        or set(alignment.attrib) != {"algn"}
        or len(alignment)
        or alignment.get("algn") != expected_alignment
    ):
        failures.append(_failure("content_mismatch", "text alignment differs", slide_id=slide_id))
    for actual, expected in zip(actual_runs, line.runs):
        text = actual.find("./a:t", NS)
        properties = actual.find("./a:rPr", NS)
        if text is None or properties is None or text.text != expected.text:
            failures.append(_failure("content_mismatch", "text content differs", slide_id=slide_id))
            continue
        if [child.tag for child in actual] != [
            "{{{}}}rPr".format(A_NS),
            "{{{}}}t".format(A_NS),
        ]:
            failures.append(_failure("content_mismatch", "text run child structure differs", slide_id=slide_id))
        allowed_attributes = {"lang", "sz", "b", "spc", "dirty"}
        if (
            not {"lang", "sz", "dirty"}.issubset(properties.attrib)
            or set(properties.attrib) - allowed_attributes
            or properties.get("lang") != "zh-CN"
            or properties.get("dirty") != "0"
            or [child.tag for child in properties]
            != [
                "{{{}}}solidFill".format(A_NS),
                "{{{}}}latin".format(A_NS),
                "{{{}}}ea".format(A_NS),
                "{{{}}}cs".format(A_NS),
            ]
        ):
            failures.append(_failure("content_mismatch", "text run formatting structure differs", slide_id=slide_id))
        expected_size = round_int((expected.style.font_size or 16.0) * 75.0)
        expected_bold = "1" if str(expected.style.font_weight) in ("bold", "600", "700", "800", "900") else None
        expected_spacing = round_int(expected.style.letter_spacing * 75.0) if expected.style.letter_spacing else None
        if (
            properties.get("sz") != str(expected_size)
            or properties.get("b") != expected_bold
            or properties.get("spc") != (None if expected_spacing is None else str(expected_spacing))
            or not _exact_solid_fill_matches(
                properties,
                (expected.style.fill or "#000000").lstrip("#").upper(),
                round_int(expected.style.fill_opacity * 100000.0),
            )
        ):
            failures.append(_failure("content_mismatch", "text run properties differ", slide_id=slide_id))
        expected_font = choose_primary_font(expected.style.font_family)
        for font_tag in ("latin", "ea", "cs"):
            font = properties.find("./a:{}".format(font_tag), NS)
            if (
                font is None
                or set(font.attrib) != {"typeface"}
                or len(font)
                or font.get("typeface") != expected_font
            ):
                failures.append(_failure("content_mismatch", "text font differs", slide_id=slide_id))
                break
        expected_space = "preserve" if expected.preserve_space or expected.text[:1].isspace() or expected.text[-1:].isspace() else None
        space_key = "{{{}}}space".format(XML_NS)
        expected_text_attributes = {space_key} if expected_space else set()
        actual_space = text.get(space_key)
        office_removed_space = bool(
            office_normalized
            and expected_space == "preserve"
            and actual_space is None
            and not text.attrib
        )
        if (
            (
                set(text.attrib) != expected_text_attributes
                or actual_space != expected_space
            )
            and not office_removed_space
        ) or len(text):
            failures.append(_failure("content_mismatch", "xml:space differs", slide_id=slide_id))
    return tuple(failures)


def verify_text_oracle(
    slide_root,
    slide_plan,
    *,
    office_normalized: bool = False,
) -> Tuple[Failure, ...]:
    expected = _expected_text_lines(slide_plan.nodes)
    failures = []
    for shape in slide_root.findall(".//p:sp", NS):
        trace = _trace(shape)
        if not trace or trace.get("kind") != "text":
            continue
        key = (trace.get("tree_path"), trace.get("line_index"))
        line = expected.pop(key, None)
        if line is None:
            failures.append(_failure("content_mismatch", "unexpected text trace", slide_id=slide_plan.slide_id))
            continue
        failures.extend(
            _verify_text_shape(
                shape,
                line,
                slide_plan.slide_id,
                office_normalized=office_normalized,
            )
        )
    if expected:
        failures.append(_failure("content_mismatch", "expected text traces are missing", slide_id=slide_plan.slide_id))
    return tuple(failures)


def _group_descendant_bounds(node: SvgNode) -> Bounds:
    retained = None
    for child in node.children:
        child_bounds = (
            _group_descendant_bounds(child)
            if child.kind == "g"
            else child.bounds
        )
        retained = (
            child_bounds
            if retained is None
            else retained.union(child_bounds)
        )
    if retained is None:
        raise ValueError("group has no source-derived descendants")
    return retained


def _expected_transform_bounds(
    nodes: Sequence[SvgNode],
) -> Dict[
    Tuple[str, str, Optional[int]],
    Tuple[Bounds, bool, bool],
]:
    result = {}
    for node in nodes:
        if node.kind == "g":
            result[(node.tree_path, "g", None)] = (
                _group_descendant_bounds(node),
                False,
                False,
            )
            result.update(_expected_transform_bounds(node.children))
        elif node.kind == "text":
            for line in node.text_lines:
                result[(node.tree_path, "text", line.line_index)] = (
                    line.bounds,
                    False,
                    False,
                )
        else:
            expansion = (
                node.style.stroke_width / 2.0
                if node.style.stroke is not None
                else 0.0
            )
            bounds = Bounds(
                node.bounds.left + expansion,
                node.bounds.top + expansion,
                node.bounds.right - expansion,
                node.bounds.bottom - expansion,
            )
            flip_h = flip_v = False
            if node.kind == "line":
                attributes = dict(node.attributes)
                flip_h = _svg_number(attributes, "x2") < _svg_number(
                    attributes,
                    "x1",
                )
                flip_v = _svg_number(attributes, "y2") < _svg_number(
                    attributes,
                    "y1",
                )
            result[(node.tree_path, node.kind, None)] = (
                bounds,
                flip_h,
                flip_v,
            )
    return result


def _shape_xfrm(shape):
    if shape.tag == "{{{}}}sp".format(P_NS):
        return shape.find("./p:spPr/a:xfrm", NS)
    return shape.find("./p:grpSpPr/a:xfrm", NS)


def verify_bounds(
    slide_root,
    slide_plan,
    config: VerificationConfig,
) -> Tuple[Failure, ...]:
    slide_id = slide_plan.slide_id
    failures = []
    tolerance = config.bounds_tolerance_px
    expected_bounds = _expected_transform_bounds(slide_plan.nodes)
    for shape in slide_root.findall(".//p:sp", NS) + slide_root.findall(".//p:grpSp", NS):
        xfrm = _shape_xfrm(shape)
        if xfrm is None:
            failures.append(_failure("bounds_violation", "shape lacks xfrm", slide_id=slide_id))
            continue
        is_group = shape.tag == "{{{}}}grpSp".format(P_NS)
        expected_children = [
            "{{{}}}off".format(A_NS),
            "{{{}}}ext".format(A_NS),
        ] + (
            ["{{{}}}chOff".format(A_NS), "{{{}}}chExt".format(A_NS)]
            if is_group
            else []
        )
        allowed_attributes = set() if is_group else {"flipH", "flipV"}
        if (
            [child.tag for child in xfrm] != expected_children
            or set(xfrm.attrib) - allowed_attributes
            or any(value != "1" for value in xfrm.attrib.values())
        ):
            failures.append(_failure("bounds_violation", "shape xfrm structure differs", slide_id=slide_id))
            continue
        off = xfrm[0]
        ext = xfrm[1]
        if set(off.attrib) != {"x", "y"} or set(ext.attrib) != {"cx", "cy"}:
            failures.append(_failure("bounds_violation", "shape xfrm attributes differ", slide_id=slide_id))
            continue
        left_emu = _required_integer(off, "x", DRAWINGML_COORD_MIN, DRAWINGML_COORD_MAX)
        top_emu = _required_integer(off, "y", DRAWINGML_COORD_MIN, DRAWINGML_COORD_MAX)
        width_emu = _required_integer(ext, "cx", 1, DRAWINGML_POSITIVE_COORD_MAX)
        height_emu = _required_integer(ext, "cy", 1, DRAWINGML_POSITIVE_COORD_MAX)
        if None in (left_emu, top_emu, width_emu, height_emu):
            failures.append(_failure("bounds_violation", "shape xfrm numeric range differs", slide_id=slide_id))
            continue
        left = left_emu / EMU_PER_PX
        top = top_emu / EMU_PER_PX
        width = width_emu / EMU_PER_PX
        height = height_emu / EMU_PER_PX
        trace = _trace(shape)
        key = (
            trace.get("tree_path") if trace else None,
            trace.get("kind") if trace else None,
            trace.get("line_index") if trace else None,
        )
        expected_oracle = expected_bounds.pop(key, None)
        if expected_oracle is None:
            failures.append(
                _failure(
                    "bounds_violation",
                    "shape has no source-derived bounds oracle",
                    slide_id=slide_id,
                    tree_path=key[0],
                    element_type=key[1],
                )
            )
        else:
            expected, expected_flip_h, expected_flip_v = expected_oracle
            actual_flip_h = xfrm.get("flipH") == "1"
            actual_flip_v = xfrm.get("flipV") == "1"
            if (
                actual_flip_h != expected_flip_h
                or actual_flip_v != expected_flip_v
            ):
                failures.append(
                    _failure(
                        "bounds_violation",
                        "shape flips differ from source-derived transform",
                        slide_id=slide_id,
                        tree_path=key[0],
                        element_type=key[1],
                    )
                )
            if any(
                difference > tolerance
                for difference in (
                    abs(left - expected.left),
                    abs(top - expected.top),
                    abs(left + width - expected.right),
                    abs(top + height - expected.bottom),
                )
            ):
                failures.append(
                    _failure(
                        "bounds_violation",
                        "shape transform differs from source-derived bounds",
                        slide_id=slide_id,
                        tree_path=key[0],
                        element_type=key[1],
                    )
                )
        if left < -tolerance or top < -tolerance or left + width > 1280 + tolerance or top + height > 720 + tolerance:
            failures.append(_failure("bounds_violation", "shape is outside slide bounds", slide_id=slide_id))
        if is_group:
            child_off = xfrm[2]
            child_ext = xfrm[3]
            if (
                set(child_off.attrib) != {"x", "y"}
                or set(child_ext.attrib) != {"cx", "cy"}
                or (off.get("x"), off.get("y")) != (child_off.get("x"), child_off.get("y"))
                or (ext.get("cx"), ext.get("cy")) != (child_ext.get("cx"), child_ext.get("cy"))
            ):
                failures.append(_failure("group_mismatch", "group identity transform differs", slide_id=slide_id))
    if expected_bounds:
        failures.append(
            _failure(
                "bounds_violation",
                "source-derived bounds are missing from candidate",
                slide_id=slide_id,
            )
        )
    return tuple(failures)


def _required_integer(node, name: str, minimum: int, maximum: int) -> Optional[int]:
    value = node.get(name)
    if (
        value is None
        or len(value) > 32
        or _INTEGER_RE.fullmatch(value) is None
    ):
        return None
    try:
        parsed = int(value)
    except ValueError:
        return None
    return parsed if minimum <= parsed <= maximum else None


def _empty_node(node) -> bool:
    return node is not None and not node.attrib and not len(node)


def _office_root_group_properties_valid(node) -> bool:
    if node is None or node.attrib or [child.tag for child in node] != [
        "{{{}}}xfrm".format(A_NS)
    ]:
        return False
    xfrm = node[0]
    expected = (
        ("off", {"x": "0", "y": "0"}),
        ("ext", {"cx": "0", "cy": "0"}),
        ("chOff", {"x": "0", "y": "0"}),
        ("chExt", {"cx": "0", "cy": "0"}),
    )
    if xfrm.attrib or len(xfrm) != len(expected):
        return False
    return all(
        child.tag == "{{{}}}{}".format(A_NS, name)
        and dict(child.attrib) == attributes
        and not len(child)
        for child, (name, attributes) in zip(xfrm, expected)
    )


def _verify_nonvisual_structures(
    sp_tree,
    slide_id: str,
    *,
    office_normalized: bool = False,
) -> Tuple[Failure, ...]:
    failures = []
    if len(sp_tree) < 2 or [sp_tree[0].tag, sp_tree[1].tag] != [
        "{{{}}}nvGrpSpPr".format(P_NS),
        "{{{}}}grpSpPr".format(P_NS),
    ]:
        return (_failure("structure_mismatch", "spTree root structure differs", slide_id=slide_id),)
    root_nv = sp_tree[0]
    if (
        root_nv.attrib
        or [child.tag for child in root_nv]
        != [
            "{{{}}}cNvPr".format(P_NS),
            "{{{}}}cNvGrpSpPr".format(P_NS),
            "{{{}}}nvPr".format(P_NS),
        ]
        or set(root_nv[0].attrib) != {"id", "name"}
        or len(root_nv[0])
        or not _empty_node(root_nv[1])
        or not _empty_node(root_nv[2])
        or not (
            _empty_node(sp_tree[1])
            or (
                office_normalized
                and _office_root_group_properties_valid(sp_tree[1])
            )
        )
    ):
        failures.append(_failure("structure_mismatch", "spTree root nonvisual metadata differs", slide_id=slide_id))
    for shape in sp_tree.findall(".//p:sp", NS):
        trace = _trace(shape)
        is_text = bool(trace and trace.get("kind") == "text")
        nv = shape.find("./p:nvSpPr", NS)
        if (
            nv is None
            or nv.attrib
            or [child.tag for child in nv]
            != [
                "{{{}}}cNvPr".format(P_NS),
                "{{{}}}cNvSpPr".format(P_NS),
                "{{{}}}nvPr".format(P_NS),
            ]
            or set(nv[0].attrib) != {"id", "name", "descr"}
            or len(nv[0])
            or set(nv[1].attrib) != ({"txBox"} if is_text else set())
            or (is_text and nv[1].get("txBox") != "1")
            or len(nv[1])
            or not _empty_node(nv[2])
        ):
            failures.append(_failure("structure_mismatch", "shape nonvisual metadata differs", slide_id=slide_id))
    for group in sp_tree.findall(".//p:grpSp", NS):
        nv = group.find("./p:nvGrpSpPr", NS)
        properties = group.find("./p:grpSpPr", NS)
        if (
            nv is None
            or properties is None
            or nv.attrib
            or [child.tag for child in nv]
            != [
                "{{{}}}cNvPr".format(P_NS),
                "{{{}}}cNvGrpSpPr".format(P_NS),
                "{{{}}}nvPr".format(P_NS),
            ]
            or set(nv[0].attrib) != {"id", "name", "descr"}
            or len(nv[0])
            or not _empty_node(nv[1])
            or not _empty_node(nv[2])
            or properties.attrib
            or [child.tag for child in properties] != ["{{{}}}xfrm".format(A_NS)]
        ):
            failures.append(_failure("structure_mismatch", "group nonvisual metadata differs", slide_id=slide_id))
    return tuple(failures)


def _expected_node_kinds(nodes: Sequence[SvgNode]) -> Mapping[str, str]:
    result = {}
    for node in nodes:
        result[node.tree_path] = node.kind
        if node.kind == "g":
            result.update(_expected_node_kinds(node.children))
    return result


def _expected_nodes_by_path(nodes: Sequence[SvgNode]) -> Mapping[str, SvgNode]:
    result = {}
    for node in nodes:
        result[node.tree_path] = node
        if node.kind == "g":
            result.update(_expected_nodes_by_path(node.children))
    return result


def _exact_no_fill_matches(parent) -> bool:
    no_fills = parent.findall("./a:noFill", NS)
    return bool(
        len(no_fills) == 1
        and not no_fills[0].attrib
        and not len(no_fills[0])
        and not parent.findall("./a:solidFill", NS)
    )


def _verify_geometry_style_oracle(slide_root, slide_plan) -> Tuple[Failure, ...]:
    expected_nodes = _expected_nodes_by_path(slide_plan.nodes)
    failures = []
    for shape in slide_root.findall(".//p:sp", NS):
        trace = _trace(shape)
        if not trace or trace.get("kind") == "text":
            continue
        tree_path = trace.get("tree_path")
        node = expected_nodes.get(tree_path)
        if node is None or node.kind != trace.get("kind"):
            continue
        properties = shape.find("./p:spPr", NS)
        line = None if properties is None else properties.find("./a:ln", NS)
        if properties is None or line is None:
            failures.append(
                _failure(
                    "content_mismatch",
                    "geometry style container is missing",
                    slide_id=slide_plan.slide_id,
                    tree_path=tree_path,
                    element_type=node.kind,
                )
            )
            continue
        fill_matches = (
            _exact_no_fill_matches(properties)
            if node.style.fill is None
            else _exact_solid_fill_matches(
                properties,
                node.style.fill.lstrip("#").upper(),
                round_int(node.style.fill_opacity * 100000.0),
            )
        )
        expected_width = str(round_int(node.style.stroke_width * EMU_PER_PX))
        stroke_fill_matches = (
            _exact_no_fill_matches(line)
            if node.style.stroke is None
            else _exact_solid_fill_matches(
                line,
                node.style.stroke.lstrip("#").upper(),
                round_int(node.style.stroke_opacity * 100000.0),
            )
        )
        if (
            not fill_matches
            or set(line.attrib) != {"w"}
            or line.get("w") != expected_width
            or len(line) != 1
            or not stroke_fill_matches
        ):
            failures.append(
                _failure(
                    "content_mismatch",
                    "geometry fill/stroke style differs from source",
                    slide_id=slide_plan.slide_id,
                    tree_path=tree_path,
                    element_type=node.kind,
                )
            )
    return tuple(failures)


def _verify_shape_direct_structure(
    slide_root,
    slide_plan,
) -> Tuple[Failure, ...]:
    slide_id = slide_plan.slide_id
    expected_kinds = _expected_node_kinds(slide_plan.nodes)
    failures = []
    for shape in slide_root.findall(".//p:sp", NS):
        trace = _trace(shape)
        is_text = bool(trace and trace.get("kind") == "text")
        expected = [
            "{{{}}}nvSpPr".format(P_NS),
            "{{{}}}spPr".format(P_NS),
        ] + (["{{{}}}txBody".format(P_NS)] if is_text else [])
        if [child.tag for child in shape] != expected:
            failures.append(
                _failure(
                    "content_mismatch" if is_text else "structure_mismatch",
                    "shape direct child structure differs",
                    slide_id=slide_id,
                )
            )
            continue
        tree_path = trace.get("tree_path") if trace else None
        expected_kind = expected_kinds.get(tree_path)
        if is_text:
            expected_kind = "text" if expected_kind == "text" else None
        if expected_kind is None:
            failures.append(_failure("structure_mismatch", "shape trace kind is unknown", slide_id=slide_id))
            continue
        properties = shape.find("./p:spPr", NS)
        if properties is None or properties.attrib or len(properties) != 4:
            failures.append(_failure("structure_mismatch", "shape properties structure differs", slide_id=slide_id))
            continue
        if properties[0].tag != "{{{}}}xfrm".format(A_NS) or properties[3].tag != "{{{}}}ln".format(A_NS):
            failures.append(_failure("structure_mismatch", "shape transform/line order differs", slide_id=slide_id))
        if properties[2].tag not in (
            "{{{}}}solidFill".format(A_NS),
            "{{{}}}noFill".format(A_NS),
        ):
            failures.append(_failure("structure_mismatch", "shape fill structure differs", slide_id=slide_id))
        custom = [child for child in properties if child.tag == "{{{}}}custGeom".format(A_NS)]
        presets = [child for child in properties if child.tag == "{{{}}}prstGeom".format(A_NS)]
        if expected_kind in ("path", "polygon", "polyline"):
            if len(custom) != 1 or presets or properties[1] is not custom[0]:
                failures.append(_failure("structure_mismatch", "custom geometry trace lacks one custGeom", slide_id=slide_id))
        else:
            expected_preset = {
                "rect": "rect",
                "circle": "ellipse",
                "ellipse": "ellipse",
                "line": "line",
                "text": "rect",
            }.get(expected_kind)
            if (
                expected_preset is None
                or custom
                or len(presets) != 1
                or properties[1] is not presets[0]
                or set(presets[0].attrib) != {"prst"}
                or presets[0].get("prst") != expected_preset
                or [child.tag for child in presets[0]] != ["{{{}}}avLst".format(A_NS)]
                or presets[0][0].attrib
                or len(presets[0][0])
            ):
                failures.append(_failure("structure_mismatch", "preset geometry trace differs", slide_id=slide_id))
    return tuple(failures)


_CUSTOM_SCALE = 100


def _svg_number(attributes: Mapping[str, str], name: str, default: float = 0.0) -> float:
    value = attributes.get(name)
    if value is None:
        return default
    lexical = value[:-2] if value.endswith("px") else value
    return float(lexical)


def _svg_points(value: str) -> Tuple[Tuple[float, float], ...]:
    tokens = [token for token in re.split(r"[\s,]+", value.strip()) if token]
    return tuple(
        (float(tokens[index]), float(tokens[index + 1]))
        for index in range(0, len(tokens), 2)
    )


def _custom_point(x: float, y: float, bounds: Bounds) -> Tuple[int, int]:
    return (
        round_int((x - bounds.left) * _CUSTOM_SCALE),
        round_int((y - bounds.top) * _CUSTOM_SCALE),
    )


def _custom_extent(value: float) -> int:
    return max(1, round_int(value * _CUSTOM_SCALE))


def _expected_custom_geometry_signature(node: SvgNode):
    attributes = dict(node.attributes)
    commands = []
    if node.kind in ("polygon", "polyline"):
        points = _svg_points(attributes["points"])
        bounds = Bounds(
            min(point[0] for point in points),
            min(point[1] for point in points),
            max(point[0] for point in points),
            max(point[1] for point in points),
        )
        commands.append(("moveTo",) + _custom_point(*points[0], bounds))
        commands.extend(
            ("lnTo",) + _custom_point(x, y, bounds)
            for x, y in points[1:]
        )
        if node.kind == "polygon":
            commands.append(("close",))
    elif node.kind == "path":
        segments = parse_path(attributes["d"])
        bounds = path_bounds(segments)
        current_x = current_y = None
        subpath_x = subpath_y = None
        for segment in segments:
            if isinstance(segment, MoveTo):
                commands.append(
                    ("moveTo",) + _custom_point(segment.x, segment.y, bounds)
                )
                current_x, current_y = segment.x, segment.y
                subpath_x, subpath_y = current_x, current_y
            elif isinstance(segment, LineTo):
                commands.append(
                    ("lnTo",) + _custom_point(segment.x, segment.y, bounds)
                )
                current_x, current_y = segment.x, segment.y
            elif isinstance(segment, ArcTo):
                center = endpoint_arc_to_center(current_x, current_y, segment)
                if center is None:
                    commands.append(
                        ("lnTo",)
                        + _custom_point(segment.end_x, segment.end_y, bounds)
                    )
                else:
                    commands.append(
                        (
                            "arcTo",
                            round_int(center.corrected_rx * _CUSTOM_SCALE),
                            round_int(center.corrected_ry * _CUSTOM_SCALE),
                            round_int(math.degrees(center.start_radians) * 60000.0),
                            round_int(math.degrees(center.sweep_radians) * 60000.0),
                        )
                    )
                    commands.append(
                        ("lnTo",)
                        + _custom_point(segment.end_x, segment.end_y, bounds)
                    )
                current_x, current_y = segment.end_x, segment.end_y
            elif isinstance(segment, ClosePath):
                commands.append(("close",))
                current_x, current_y = subpath_x, subpath_y
    else:
        return None
    width = _custom_extent(bounds.right - bounds.left)
    height = _custom_extent(bounds.bottom - bounds.top)
    return ((0, 0, width, height), (width, height), tuple(commands))


def _actual_custom_geometry_signature(geometry):
    rectangle = geometry.find("./a:rect", NS)
    path = geometry.find("./a:pathLst/a:path", NS)
    if rectangle is None or path is None:
        return None
    rect_values = tuple(
        _required_integer(
            rectangle,
            name,
            CUSTOM_COORD_MIN,
            CUSTOM_COORD_MAX,
        )
        for name in ("l", "t", "r", "b")
    )
    width = _required_integer(path, "w", 1, CUSTOM_COORD_MAX)
    height = _required_integer(path, "h", 1, CUSTOM_COORD_MAX)
    if None in rect_values or width is None or height is None:
        return None
    commands = []
    for command in path:
        local = command.tag.rsplit("}", 1)[-1]
        if local in ("moveTo", "lnTo"):
            point = command.find("./a:pt", NS)
            if point is None:
                return None
            x = _required_integer(point, "x", CUSTOM_COORD_MIN, CUSTOM_COORD_MAX)
            y = _required_integer(point, "y", CUSTOM_COORD_MIN, CUSTOM_COORD_MAX)
            if x is None or y is None:
                return None
            commands.append((local, x, y))
        elif local == "arcTo":
            values = tuple(
                _required_integer(command, name, minimum, maximum)
                for name, minimum, maximum in (
                    ("wR", 1, CUSTOM_COORD_MAX),
                    ("hR", 1, CUSTOM_COORD_MAX),
                    ("stAng", DRAWINGML_ANGLE_MIN, DRAWINGML_ANGLE_MAX),
                    ("swAng", DRAWINGML_ANGLE_MIN, DRAWINGML_ANGLE_MAX),
                )
            )
            if None in values:
                return None
            commands.append((local,) + values)
        elif local == "close":
            commands.append((local,))
        else:
            return None
    return (rect_values, (width, height), tuple(commands))


def _has_nonwhitespace_xml_text(node) -> bool:
    return bool((node.text or "").strip() or (node.tail or "").strip())


def _verify_custom_geometry(slide_root, slide_plan) -> Tuple[Failure, ...]:
    slide_id = slide_plan.slide_id
    expected_nodes = _expected_nodes_by_path(slide_plan.nodes)
    failures = []
    expected_geometry_children = [
        "{{{}}}avLst".format(A_NS),
        "{{{}}}gdLst".format(A_NS),
        "{{{}}}ahLst".format(A_NS),
        "{{{}}}cxnLst".format(A_NS),
        "{{{}}}rect".format(A_NS),
        "{{{}}}pathLst".format(A_NS),
    ]
    for geometry in slide_root.findall(".//a:custGeom", NS):
        if geometry.attrib:
            failures.append(_failure("structure_mismatch", "custom geometry attributes are forbidden", slide_id=slide_id))
        if any(_has_nonwhitespace_xml_text(node) for node in geometry.iter()):
            failures.append(
                _failure(
                    "structure_mismatch",
                    "custom geometry contains unexpected text",
                    slide_id=slide_id,
                )
            )
        if [child.tag for child in geometry] != expected_geometry_children:
            failures.append(_failure("structure_mismatch", "custom geometry child order differs", slide_id=slide_id))
            continue
        if any(child.attrib or len(child) for child in list(geometry)[:4]):
            failures.append(_failure("structure_mismatch", "custom geometry list elements must be empty", slide_id=slide_id))
        rectangle = geometry.find("./a:rect", NS)
        if (
            rectangle is None
            or len(rectangle)
            or set(rectangle.attrib) != {"l", "t", "r", "b"}
            or any(
                _required_integer(
                    rectangle,
                    name,
                    CUSTOM_COORD_MIN,
                    CUSTOM_COORD_MAX,
                )
                is None
                for name in ("l", "t", "r", "b")
            )
        ):
            failures.append(_failure("structure_mismatch", "custom geometry rect is invalid", slide_id=slide_id))
        path_list = geometry.find("./a:pathLst", NS)
        paths = [] if path_list is None else path_list.findall("./a:path", NS)
        if (
            path_list is None
            or path_list.attrib
            or len(path_list) != 1
            or len(paths) != 1
        ):
            failures.append(_failure("structure_mismatch", "custom geometry requires exactly one path", slide_id=slide_id))
            continue
        path = paths[0]
        width = _required_integer(path, "w", 1, CUSTOM_COORD_MAX)
        height = _required_integer(path, "h", 1, CUSTOM_COORD_MAX)
        if set(path.attrib) != {"w", "h"} or width is None or height is None:
            failures.append(_failure("structure_mismatch", "custom path extents are invalid", slide_id=slide_id))
        commands = list(path)
        allowed_command_tags = {
            "{{{}}}moveTo".format(A_NS),
            "{{{}}}lnTo".format(A_NS),
            "{{{}}}arcTo".format(A_NS),
            "{{{}}}close".format(A_NS),
        }
        if (
            not commands
            or commands[0].tag != "{{{}}}moveTo".format(A_NS)
            or any(command.tag not in allowed_command_tags for command in commands)
            or not any(
                command.tag in (
                    "{{{}}}lnTo".format(A_NS),
                    "{{{}}}arcTo".format(A_NS),
                )
                for command in commands
            )
        ):
            failures.append(_failure("structure_mismatch", "custom path has no drawable commands", slide_id=slide_id))
        for command in commands:
            local = command.tag.rsplit("}", 1)[-1]
            if local in ("moveTo", "lnTo"):
                points = command.findall("./a:pt", NS)
                if (
                    set(command.attrib)
                    or len(command) != 1
                    or len(points) != 1
                    or len(points[0])
                    or set(points[0].attrib) != {"x", "y"}
                    or _required_integer(points[0], "x", CUSTOM_COORD_MIN, CUSTOM_COORD_MAX) is None
                    or _required_integer(points[0], "y", CUSTOM_COORD_MIN, CUSTOM_COORD_MAX) is None
                ):
                    failures.append(_failure("structure_mismatch", "custom point command is invalid", slide_id=slide_id))
            elif local == "arcTo":
                required = {"wR", "hR", "stAng", "swAng"}
                radius_x = _required_integer(command, "wR", 1, CUSTOM_COORD_MAX)
                radius_y = _required_integer(command, "hR", 1, CUSTOM_COORD_MAX)
                start = _required_integer(command, "stAng", DRAWINGML_ANGLE_MIN, DRAWINGML_ANGLE_MAX)
                sweep = _required_integer(command, "swAng", DRAWINGML_ANGLE_MIN, DRAWINGML_ANGLE_MAX)
                if (
                    set(command.attrib) != required
                    or radius_x is None
                    or radius_y is None
                    or start is None
                    or sweep is None
                    or sweep == 0
                    or len(command)
                ):
                    failures.append(_failure("structure_mismatch", "custom arc command is invalid", slide_id=slide_id))
            elif local == "close":
                if command.attrib or len(command):
                    failures.append(_failure("structure_mismatch", "custom close command is invalid", slide_id=slide_id))
            else:
                failures.append(_failure("structure_mismatch", "unsupported custom path command", slide_id=slide_id))
        shape = geometry.getparent().getparent()
        trace = _trace(shape)
        tree_path = trace.get("tree_path") if trace else None
        expected_node = expected_nodes.get(tree_path)
        if expected_node is None or expected_node.kind not in (
            "path",
            "polygon",
            "polyline",
        ):
            failures.append(
                _failure(
                    "structure_mismatch",
                    "custom geometry has no source-derived oracle",
                    slide_id=slide_id,
                    tree_path=tree_path,
                )
            )
            continue
        try:
            expected_signature = _expected_custom_geometry_signature(expected_node)
        except (KeyError, TypeError, ValueError, EditableError):
            expected_signature = None
        actual_signature = _actual_custom_geometry_signature(geometry)
        if (
            expected_signature is None
            or actual_signature != expected_signature
        ):
            failures.append(
                _failure(
                    "structure_mismatch",
                    "custom geometry differs from source-derived oracle",
                    slide_id=slide_id,
                    tree_path=tree_path,
                    element_type=expected_node.kind,
                )
            )
    return tuple(failures)


def verify_slide_tree(
    slide_root,
    slide_plan,
    *,
    office_normalized: bool = False,
):
    failures = []
    if slide_root.tag != "{{{}}}sld".format(P_NS):
        failures.append(
            _failure(
                "structure_mismatch",
                "slide root element differs",
                slide_id=slide_plan.slide_id,
            )
        )
        return tuple(failures), 0, 0, 0
    if slide_root.findall(".//p:pic", NS) or slide_root.findall(".//a:blip", NS):
        failures.append(
            _failure(
                "image_fallback_detected",
                "slide contains an image shape",
                slide_id=slide_plan.slide_id,
            )
        )
    content_trees = slide_root.findall("./p:cSld", NS)
    if len(content_trees) != 1:
        failures.append(
            _failure(
                "structure_mismatch",
                "slide must contain exactly one cSld",
                slide_id=slide_plan.slide_id,
            )
        )
        if not content_trees:
            return tuple(failures), 0, 0, 0
    c_sld = content_trees[0]
    if c_sld.get("name") != (slide_plan.title or slide_plan.slide_id):
        failures.append(_failure("content_mismatch", "slide title metadata differs", slide_id=slide_plan.slide_id))
        return tuple(failures), 0, 0, 0
    sp_tree = c_sld.find("./p:spTree", NS)
    if sp_tree is None:
        failures.append(_failure("structure_mismatch", "slide lacks spTree", slide_id=slide_plan.slide_id))
        return tuple(failures), 0, 0, 0
    failures.extend(
        _verify_nonvisual_structures(
            sp_tree,
            slide_plan.slide_id,
            office_normalized=office_normalized,
        )
    )
    root_identity = sp_tree.find("./p:nvGrpSpPr/p:cNvPr", NS)
    if root_identity is None or root_identity.get("id") != "1":
        failures.append(_failure("structure_mismatch", "spTree root ID must equal 1", slide_id=slide_plan.slide_id))
    ids = [node.get("id") for node in sp_tree.findall(".//p:cNvPr", NS)]
    if (
        None in ids
        or len(ids) != len(set(ids))
        or any(re.fullmatch(r"[0-9]+", value or "") is None for value in ids)
        or any(
            len(value) > 10 or int(value) > 4294967295
            for value in ids
            if value is not None and value.isdigit()
        )
    ):
        failures.append(_failure("structure_mismatch", "shape IDs are invalid or duplicated", slide_id=slide_plan.slide_id))
    counters = {"leaves": 0, "groups": 0, "pics": 0}
    actual = _actual_tree(_shape_children(sp_tree), counters)
    expected = _expected_tree(slide_plan.slide_id, slide_plan.nodes)
    if _identity_only(actual) != _identity_only(expected):
        failures.append(_failure("structure_mismatch", "shape ID/name/trace sequence differs", slide_id=slide_plan.slide_id))
        failures.append(_failure("group_mismatch", "trace hierarchy/order differs", slide_id=slide_plan.slide_id))
    for item in actual:
        if "children" in item and not item["children"]:
            failures.append(_failure("group_mismatch", "empty group emitted", slide_id=slide_plan.slide_id))
    failures.extend(_verify_shape_direct_structure(slide_root, slide_plan))
    failures.extend(_verify_geometry_style_oracle(slide_root, slide_plan))
    failures.extend(_verify_custom_geometry(slide_root, slide_plan))
    return tuple(failures), counters["leaves"], counters["groups"], len(_shape_children(sp_tree))


def verify_notes(presentation: Presentation, deck_plan: DeckPlan) -> Tuple[Failure, ...]:
    failures = []
    for slide, plan in zip(presentation.slides, deck_plan.slides):
        expected = tuple(line for line in format_speaker_notes(plan.notes).splitlines() if line)
        actual = tuple(
            paragraph.text
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs
            if paragraph.text
        )
        if actual != expected:
            failures.append(_failure("notes_mismatch", "speaker notes differ", slide_id=plan.slide_id))
    return tuple(failures)


def verify_candidate(
    pptx_path: Path,
    deck_plan: DeckPlan,
    config: VerificationConfig,
    *,
    office_normalized: bool = False,
) -> VerificationReport:
    slide_metadata = tuple(
        (slide.slide_id, slide.title, slide.description)
        for slide in deck_plan.slides
    )
    payload, zip_failures = _zip_payload(Path(pptx_path))
    if payload is None or any(failure.code == "pptx_zip_invalid" for failure in zip_failures):
        return VerificationReport(
            False,
            tuple(zip_failures),
            0,
            0,
            0,
            0,
            slide_metadata,
        )
    failures = list(
        _verify_payload_parts(
            payload,
            zip_failures,
            office_normalized=office_normalized,
        )
    )
    try:
        presentation = Presentation(str(pptx_path))
    except Exception:
        failures.append(_failure("pptx_reopen_failed", "python-pptx cannot reopen candidate"))
        return VerificationReport(
            False,
            tuple(failures),
            0,
            0,
            0,
            0,
            slide_metadata,
        )
    if len(presentation.slides) != len(deck_plan.slides):
        failures.append(_failure("structure_mismatch", "slide count differs"))
    if (
        presentation.slide_width,
        presentation.slide_height,
    ) != (12192000, 6858000):
        failures.append(
            _failure(
                "structure_mismatch",
                "slide dimensions differ from exact 12192000x6858000 EMU",
            )
        )
    leaf_count = group_count = top_count = 0
    ordered_parts, order_failures = _ordered_slide_parts(payload)
    failures.extend(order_failures)
    if len(ordered_parts) != len(deck_plan.slides):
        failures.append(_failure("structure_mismatch", "ordered slide-part count differs"))
    for index, plan in enumerate(deck_plan.slides, 1):
        member = (
            ordered_parts[index - 1]
            if index <= len(ordered_parts)
            else "ppt/slides/slide{}.xml".format(index)
        )
        canonical_member = "ppt/slides/slide{}.xml".format(index)
        if member != canonical_member:
            failures.append(
                _failure(
                    "structure_mismatch",
                    "displayed slide order differs from canonical plan order",
                    slide_id=plan.slide_id,
                )
            )
        data = payload.get(member)
        if data is None:
            failures.append(_failure("structure_mismatch", "slide part is missing", slide_id=plan.slide_id))
            continue
        try:
            root = _xml(data)
        except etree.XMLSyntaxError:
            failures.append(_failure("structure_mismatch", "slide XML is malformed", slide_id=plan.slide_id))
            continue
        try:
            slide_failures, leaves, groups, top = verify_slide_tree(
                root,
                plan,
                office_normalized=office_normalized,
            )
            failures.extend(slide_failures)
            failures.extend(
                verify_text_oracle(
                    root,
                    plan,
                    office_normalized=office_normalized,
                )
            )
            failures.extend(verify_bounds(root, plan, config))
        except Exception:
            failures.append(
                _failure(
                    "structure_mismatch",
                    "candidate slide contains an unprocessable value",
                    slide_id=plan.slide_id,
                )
            )
            continue
        leaf_count += leaves
        group_count += groups
        top_count += top
    try:
        failures.extend(verify_notes(presentation, deck_plan))
    except Exception:
        failures.append(
            _failure(
                "pptx_reopen_failed",
                "python-pptx cannot traverse candidate slides or notes",
            )
        )
    return VerificationReport(
        passed=not failures,
        failures=tuple(failures),
        slide_count=len(presentation.slides),
        top_level_shape_count=top_count,
        recursive_leaf_count=leaf_count,
        recursive_group_count=group_count,
        slide_metadata=slide_metadata,
    )
