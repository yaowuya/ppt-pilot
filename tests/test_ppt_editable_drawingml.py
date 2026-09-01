import io
import json
import os
import re
import shutil
import sys
import tempfile
import unittest
from unittest import mock
import zipfile
from pathlib import Path


REPO_ROOT = Path(__file__).resolve().parents[1]
SKILL_ROOT = REPO_ROOT / "skills" / "ppt-editable"
SCRIPTS_ROOT = SKILL_ROOT / "scripts"
FIXTURE_ROOT = REPO_ROOT / "tests" / "fixtures" / "ppt-editable" / "svg"
RUN_FIXTURE_ROOT = REPO_ROOT / "tests" / "fixtures" / "ppt-editable" / "run-complete"
sys.path.insert(0, str(SCRIPTS_ROOT))

from _ppt_editable import Bounds, DeckPlan, ResolvedStyle, SlidePlan, SpeakerNotes, SvgNode  # noqa: E402
from _ppt_editable.contract import SlideSource  # noqa: E402
from _ppt_editable.drawingml import (  # noqa: E402
    EMU_PER_PX,
    ShapeIdAllocator,
    build_slide,
    stable_shape_name,
    trace_description,
)
from _ppt_editable.svg_parser import parse_svg_slide  # noqa: E402


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
XML_NS = "http://www.w3.org/XML/1998/namespace"


def qn(namespace, local):
    return "{{{}}}{}".format(namespace, local)


def descendants(element, namespace, local):
    return element.findall(".//" + qn(namespace, local))


class DrawingMlTests(unittest.TestCase):
    def _plan(self, filename):
        path = FIXTURE_ROOT / filename
        return parse_svg_slide(
            SlideSource("S01", path, "slides/S01.svg", "production"),
            SpeakerNotes(),
        )

    def _elements(self, filename, include_text=True):
        return build_slide(
            self._plan(filename),
            ShapeIdAllocator(),
            include_text=include_text,
        )

    def _elements_from_svg(self, body):
        directory = tempfile.TemporaryDirectory()
        self.addCleanup(directory.cleanup)
        path = Path(directory.name) / "S01.svg"
        path.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
            '<title>x</title><desc>x</desc>{}</svg>'.format(body),
            encoding="utf-8",
        )
        plan = parse_svg_slide(
            SlideSource("S01", path, "slides/S01.svg", "production"),
            SpeakerNotes(),
        )
        return build_slide(plan, ShapeIdAllocator())

    def test_shape_id_allocator_is_monotonic_and_positive(self):
        allocator = ShapeIdAllocator(first_id=2)
        self.assertEqual([allocator.next() for _ in range(4)], [2, 3, 4, 5])
        with self.assertRaises(ValueError):
            ShapeIdAllocator(first_id=0)
        maximum = ShapeIdAllocator(first_id=4294967295)
        self.assertEqual(maximum.next(), 4294967295)
        with self.assertRaises(ValueError):
            maximum.next()

    def test_stable_names_and_trace_are_deterministic_unique_and_sanitized(self):
        first = stable_shape_name("S01", "/svg[1]/g[1]/path[1]", "path", "SRC/01")
        second = stable_shape_name("S01", "/svg[1]/g[1]/path[2]", "path", "SRC/01")
        self.assertEqual(first, stable_shape_name("S01", "/svg[1]/g[1]/path[1]", "path", "SRC/01"))
        self.assertNotEqual(first, second)
        self.assertRegex(first, r"^SRC_01__path__[0-9a-f]{8}$")
        trace = json.loads(
            trace_description(
                "S01",
                "/svg[1]/text[1]",
                "text",
                "SRC/01",
                line_index=2,
            )
        )
        self.assertEqual(
            trace,
            {
                "schema_version": 1,
                "slide_id": "S01",
                "tree_path": "/svg[1]/text[1]",
                "kind": "text",
                "source_id": "SRC/01",
                "line_index": 2,
            },
        )

    def test_primitives_map_to_native_shapes_with_fill_alpha_and_line(self):
        elements = self._elements("primitives.svg")
        self.assertEqual(len(elements), 6)
        self.assertTrue(all(element.tag == qn(P_NS, "sp") for element in elements))
        rect = elements[0]
        xfrm = descendants(rect, A_NS, "xfrm")[0]
        off = xfrm.find(qn(A_NS, "off"))
        ext = xfrm.find(qn(A_NS, "ext"))
        self.assertEqual((off.get("x"), off.get("y")), (str(10 * EMU_PER_PX), str(20 * EMU_PER_PX)))
        self.assertEqual((ext.get("cx"), ext.get("cy")), (str(100 * EMU_PER_PX), str(50 * EMU_PER_PX)))
        self.assertEqual(descendants(rect, A_NS, "prstGeom")[0].get("prst"), "rect")
        self.assertEqual(descendants(rect, A_NS, "alpha")[0].get("val"), "50000")
        self.assertEqual(descendants(elements[1], A_NS, "prstGeom")[0].get("prst"), "ellipse")
        self.assertEqual(descendants(elements[2], A_NS, "prstGeom")[0].get("prst"), "ellipse")
        self.assertEqual(descendants(elements[3], A_NS, "prstGeom")[0].get("prst"), "line")
        self.assertTrue(descendants(elements[4], A_NS, "custGeom"))
        self.assertTrue(descendants(elements[4], A_NS, "close"))
        self.assertTrue(descendants(elements[5], A_NS, "custGeom"))
        self.assertFalse(descendants(elements[5], A_NS, "close"))
        for element in elements:
            self.assertFalse(descendants(element, P_NS, "pic"))
            self.assertFalse(descendants(element, A_NS, "blip"))

    def test_path_geometry_contains_path_list_corrected_arcs_and_integer_lexemes(self):
        elements = self._elements("arcs.svg")
        self.assertEqual(len(elements), 2)
        for element in elements:
            self.assertEqual(len(descendants(element, A_NS, "pathLst")), 1)
            self.assertEqual(len(descendants(element, A_NS, "path")), 1)
            self.assertFalse(descendants(element, P_NS, "pic"))
            for node in element.iter():
                for key in ("x", "y", "w", "h", "wR", "hR", "stAng", "swAng", "l", "t", "r", "b"):
                    value = node.get(key)
                    if value is not None:
                        self.assertRegex(value, r"^-?[0-9]+$")
        first_arcs = descendants(elements[0], A_NS, "arcTo")
        self.assertEqual(len(first_arcs), 4)
        self.assertTrue(all(arc.get("wR") == "2000" for arc in first_arcs))
        corrected = descendants(elements[1], A_NS, "arcTo")[0]
        self.assertEqual((corrected.get("wR"), corrected.get("hR")), ("5000", "5000"))
        self.assertGreaterEqual(len(descendants(elements[0], A_NS, "lnTo")), 4)

    def test_text_lines_are_true_text_boxes_with_runs_fonts_and_space(self):
        elements = self._elements("namespace-text.svg")
        self.assertEqual(len(elements), 6)
        first = elements[0]
        c_nv_sp = descendants(first, P_NS, "cNvSpPr")[0]
        self.assertEqual(c_nv_sp.get("txBox"), "1")
        body_pr = descendants(first, A_NS, "bodyPr")[0]
        self.assertEqual(
            (body_pr.get("wrap"), body_pr.get("lIns"), body_pr.get("tIns"), body_pr.get("rIns"), body_pr.get("bIns"), body_pr.get("anchor")),
            ("none", "0", "0", "0", "0", "t"),
        )
        self.assertEqual(len(descendants(first, A_NS, "noAutofit")), 1)
        texts = descendants(first, A_NS, "t")
        self.assertEqual([text.text for text in texts], ["广东电网 ", "27 套", " 系统收敛"])
        self.assertEqual(texts[0].get(qn(XML_NS, "space")), "preserve")
        run_properties = descendants(first, A_NS, "rPr")
        self.assertEqual(run_properties[0].get("sz"), "1500")
        self.assertEqual(run_properties[1].get("b"), "1")
        self.assertEqual(descendants(run_properties[0], A_NS, "latin")[0].get("typeface"), "Arial")
        middle_line = elements[4]
        self.assertEqual(descendants(middle_line, A_NS, "pPr")[0].get("algn"), "ctr")

    def test_recursive_groups_preserve_hierarchy_order_and_identity_coordinates(self):
        elements = self._elements("nested-groups.svg")
        self.assertEqual(len(elements), 1)
        outer = elements[0]
        self.assertEqual(outer.tag, qn(P_NS, "grpSp"))
        groups = [outer] + descendants(outer, P_NS, "grpSp")
        self.assertEqual(len(groups), 2)
        leaves = descendants(outer, P_NS, "sp")
        self.assertEqual(len(leaves), 3)
        ids = [int(node.get("id")) for node in descendants(outer, P_NS, "cNvPr")]
        self.assertEqual(ids, list(range(2, 2 + len(ids))))
        direct_tags = [child.tag for child in outer if child.tag in (qn(P_NS, "sp"), qn(P_NS, "grpSp"))]
        self.assertEqual(direct_tags, [qn(P_NS, "sp"), qn(P_NS, "grpSp")])
        for group in groups:
            xfrm = descendants(group, A_NS, "xfrm")[0]
            off = xfrm.find(qn(A_NS, "off"))
            ext = xfrm.find(qn(A_NS, "ext"))
            child_off = xfrm.find(qn(A_NS, "chOff"))
            child_ext = xfrm.find(qn(A_NS, "chExt"))
            self.assertEqual((off.get("x"), off.get("y")), (child_off.get("x"), child_off.get("y")))
            self.assertEqual((ext.get("cx"), ext.get("cy")), (child_ext.get("cx"), child_ext.get("cy")))
            self.assertGreater(int(ext.get("cx")), 0)
            self.assertGreater(int(ext.get("cy")), 0)
        traces = [json.loads(node.get("descr")) for node in descendants(outer, P_NS, "cNvPr")]
        self.assertEqual(len({(trace["tree_path"], trace["kind"], trace["line_index"]) for trace in traces}), len(traces))
        self.assertIn("SRC-001", {trace["source_id"] for trace in traces})

    def test_group_transform_uses_retained_descendant_union(self):
        style = ResolvedStyle(fill="#112233", stroke=None)
        first = SvgNode(
            kind="rect",
            tree_path="/svg[1]/g[1]/rect[1]",
            style=style,
            bounds=Bounds(10.0, 20.0, 30.0, 40.0),
            attributes=(("height", "20"), ("width", "20"), ("x", "10"), ("y", "20")),
        )
        second = SvgNode(
            kind="rect",
            tree_path="/svg[1]/g[1]/rect[2]",
            style=style,
            bounds=Bounds(50.0, 60.0, 80.0, 90.0),
            attributes=(("height", "30"), ("width", "30"), ("x", "50"), ("y", "60")),
        )
        group_node = SvgNode(
            kind="g",
            tree_path="/svg[1]/g[1]",
            style=style,
            bounds=Bounds(0.0, 0.0, 1200.0, 700.0),
            children=(first, second),
        )
        slide = SlidePlan(
            slide_id="S01",
            source_path=Path("S01.svg"),
            title="x",
            description="x",
            nodes=(group_node,),
            notes=SpeakerNotes(),
        )
        group = build_slide(slide, ShapeIdAllocator())[0]
        xfrm = group.find(qn(P_NS, "grpSpPr")).find(qn(A_NS, "xfrm"))
        off = xfrm.find(qn(A_NS, "off"))
        ext = xfrm.find(qn(A_NS, "ext"))
        child_off = xfrm.find(qn(A_NS, "chOff"))
        child_ext = xfrm.find(qn(A_NS, "chExt"))
        expected_off = (str(10 * EMU_PER_PX), str(20 * EMU_PER_PX))
        expected_ext = (str(70 * EMU_PER_PX), str(70 * EMU_PER_PX))
        self.assertEqual((off.get("x"), off.get("y")), expected_off)
        self.assertEqual((ext.get("cx"), ext.get("cy")), expected_ext)
        self.assertEqual((child_off.get("x"), child_off.get("y")), expected_off)
        self.assertEqual((child_ext.get("cx"), child_ext.get("cy")), expected_ext)

    def test_group_rejects_nonpositive_retained_descendant_extent(self):
        style = ResolvedStyle(fill=None, stroke="#112233", stroke_width=0.0)
        line = SvgNode(
            kind="line",
            tree_path="/svg[1]/g[1]/line[1]",
            style=style,
            bounds=Bounds(0.0, 0.0, 10.0, 0.0),
            attributes=(("x1", "0"), ("x2", "10"), ("y1", "0"), ("y2", "0")),
        )
        group_node = SvgNode(
            kind="g",
            tree_path="/svg[1]/g[1]",
            style=style,
            bounds=line.bounds,
            children=(line,),
        )
        slide = SlidePlan(
            slide_id="S01",
            source_path=Path("S01.svg"),
            title="x",
            description="x",
            nodes=(group_node,),
            notes=SpeakerNotes(),
        )
        with self.assertRaisesRegex(ValueError, "group bounds must have positive extents"):
            build_slide(slide, ShapeIdAllocator())

    def test_geometry_only_omitted_text_group_does_not_consume_shape_id(self):
        directory = tempfile.TemporaryDirectory()
        self.addCleanup(directory.cleanup)
        path = Path(directory.name) / "S01.svg"
        path.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
            '<title>x</title><desc>x</desc>'
            '<g><text x="10" y="20">text only</text></g>'
            '<rect x="30" y="40" width="10" height="10"/>'
            '</svg>',
            encoding="utf-8",
        )
        plan = parse_svg_slide(
            SlideSource("S01", path, "slides/S01.svg", "production"),
            SpeakerNotes(),
        )
        geometry = build_slide(plan, ShapeIdAllocator(), include_text=False)
        self.assertEqual(len(geometry), 1)
        self.assertEqual(geometry[0].tag, qn(P_NS, "sp"))
        self.assertEqual(descendants(geometry[0], P_NS, "cNvPr")[0].get("id"), "2")

    def test_production_empty_group_fails_closed_instead_of_disappearing(self):
        group_node = SvgNode(
            kind="g",
            tree_path="/svg[1]/g[1]",
            style=ResolvedStyle(),
            bounds=Bounds(0.0, 0.0, 10.0, 10.0),
            children=(),
        )
        slide = SlidePlan(
            slide_id="S01",
            source_path=Path("S01.svg"),
            title="x",
            description="x",
            nodes=(group_node,),
            notes=SpeakerNotes(),
        )
        with self.assertRaisesRegex(ValueError, "production group has no descendants"):
            build_slide(slide, ShapeIdAllocator(), include_text=True)

    def test_exact_direct_child_order_and_cardinality(self):
        primitive = self._elements("primitives.svg")[0]
        self.assertEqual(
            [child.tag for child in primitive],
            [qn(P_NS, "nvSpPr"), qn(P_NS, "spPr")],
        )
        text = self._elements("namespace-text.svg")[0]
        self.assertEqual(
            [child.tag for child in text],
            [qn(P_NS, "nvSpPr"), qn(P_NS, "spPr"), qn(P_NS, "txBody")],
        )
        text_body = text.find(qn(P_NS, "txBody"))
        self.assertEqual(
            [child.tag for child in text_body],
            [qn(A_NS, "bodyPr"), qn(A_NS, "lstStyle"), qn(A_NS, "p")],
        )
        self.assertEqual(
            [child.tag for child in text_body.find(qn(A_NS, "bodyPr"))],
            [qn(A_NS, "noAutofit")],
        )
        paragraph = text_body.find(qn(A_NS, "p"))
        self.assertEqual(paragraph[0].tag, qn(A_NS, "pPr"))
        self.assertTrue(all(child.tag == qn(A_NS, "r") for child in paragraph[1:]))
        group = self._elements("nested-groups.svg")[0]
        self.assertEqual(
            [child.tag for child in group[:2]],
            [qn(P_NS, "nvGrpSpPr"), qn(P_NS, "grpSpPr")],
        )
        custom = descendants(self._elements("arcs.svg")[0], A_NS, "custGeom")[0]
        self.assertEqual(
            [child.tag.rsplit("}", 1)[-1] for child in custom],
            ["avLst", "gdLst", "ahLst", "cxnLst", "rect", "pathLst"],
        )
        path_list = custom.find(qn(A_NS, "pathLst"))
        self.assertEqual(len(path_list), 1)
        self.assertEqual(path_list[0].tag, qn(A_NS, "path"))

    def test_known_arc_angles_command_order_and_endpoint_are_exact(self):
        first, corrected_shape = self._elements("arcs.svg")
        path = descendants(first, A_NS, "path")[0]
        children = list(path)
        first_arc_index = next(
            index for index, child in enumerate(children) if child.tag == qn(A_NS, "arcTo")
        )
        first_arc = children[first_arc_index]
        self.assertEqual(
            (first_arc.get("stAng"), first_arc.get("swAng")),
            ("-5400000", "5400000"),
        )
        endpoint_line = children[first_arc_index + 1]
        self.assertEqual(endpoint_line.tag, qn(A_NS, "lnTo"))
        endpoint = endpoint_line.find(qn(A_NS, "pt"))
        self.assertEqual((endpoint.get("x"), endpoint.get("y")), ("22000", "2000"))

        corrected = descendants(corrected_shape, A_NS, "arcTo")[0]
        self.assertEqual(
            (corrected.get("stAng"), corrected.get("swAng")),
            ("10800000", "10800000"),
        )
        corrected_path = descendants(corrected_shape, A_NS, "path")[0]
        corrected_children = list(corrected_path)
        arc_index = corrected_children.index(corrected)
        corrected_endpoint = corrected_children[arc_index + 1].find(qn(A_NS, "pt"))
        self.assertEqual(
            (corrected_endpoint.get("x"), corrected_endpoint.get("y")),
            ("10000", "5000"),
        )

    def test_fill_and_stroke_alpha_are_independent_and_bounded(self):
        element = self._elements_from_svg(
            '<rect x="0" y="0" width="10" height="10" fill="#112233" '
            'stroke="#445566" stroke-width="2" fill-opacity="0.4" '
            'stroke-opacity="0.6" opacity="0.5"/>'
        )[0]
        sp_pr = element.find(qn(P_NS, "spPr"))
        direct_fill = sp_pr.find(qn(A_NS, "solidFill"))
        line = sp_pr.find(qn(A_NS, "ln"))
        self.assertEqual(descendants(direct_fill, A_NS, "alpha")[0].get("val"), "20000")
        self.assertEqual(descendants(line, A_NS, "alpha")[0].get("val"), "30000")
        opaque = self._elements_from_svg(
            '<rect x="0" y="0" width="10" height="10" fill="#112233"/>'
        )[0]
        self.assertFalse(descendants(opaque, A_NS, "alpha"))
        transparent = self._elements_from_svg(
            '<rect x="0" y="0" width="10" height="10" fill="#112233" opacity="0"/>'
        )[0]
        self.assertEqual(descendants(transparent, A_NS, "alpha")[0].get("val"), "0")
        stroke_zero = self._elements_from_svg(
            '<line x1="0" y1="0" x2="10" y2="10" stroke="#112233" '
            'stroke-width="1" stroke-opacity="0"/>'
        )[0]
        self.assertEqual(descendants(stroke_zero, A_NS, "alpha")[0].get("val"), "0")
        stroke_opaque = self._elements_from_svg(
            '<line x1="0" y1="0" x2="10" y2="10" stroke="#112233" '
            'stroke-width="1" stroke-opacity="1"/>'
        )[0]
        self.assertFalse(descendants(stroke_opaque, A_NS, "alpha"))

    def test_out_of_schema_numeric_values_are_rejected_before_serialization(self):
        cases = (
            '<text x="0" y="20" font-size="0.1">tiny</text>',
            '<text x="0" y="20" font-size="20" letter-spacing="10000000">wide</text>',
            '<line x1="0" y1="0" x2="10" y2="10" stroke="#112233" stroke-width="10000000"/>',
            '<rect x="1e20" y="0" width="10" height="10"/>',
            '<path d="M0 0A300000000000 300000000000 0 0 1 1 0" fill="none" stroke="#112233"/>',
            '<path d="M0 0A100000000000 100000000000 0 0 1 1 0" fill="none" stroke="#112233"/>',
        )
        for body in cases:
            with self.subTest(body=body):
                with self.assertRaises(ValueError):
                    self._elements_from_svg(body)

    @unittest.skipUnless(os.environ.get("PPT_EDITABLE_REFERENCE_RUN"), "reference run not configured")
    def test_reference_corpus_has_570_visible_leaves_63_groups_and_no_images(self):
        from _ppt_editable.contract import (
            parse_storyboard,
            resolve_slide_sources,
            validate_completed_run,
        )
        from _ppt_editable.svg_parser import preflight_deck

        run = Path(os.environ["PPT_EDITABLE_REFERENCE_RUN"])
        context = validate_completed_run(run)
        storyboard = parse_storyboard(context.storyboard_path)
        sources = resolve_slide_sources(context, storyboard)
        deck = preflight_deck(context, sources, storyboard, "sha256:" + "a" * 64)
        leaves = groups = pics = blips = 0
        for slide in deck.slides:
            for root in build_slide(slide, ShapeIdAllocator()):
                leaves += (root.tag == qn(P_NS, "sp")) + len(descendants(root, P_NS, "sp"))
                groups += (root.tag == qn(P_NS, "grpSp")) + len(descendants(root, P_NS, "grpSp"))
                pics += (root.tag == qn(P_NS, "pic")) + len(descendants(root, P_NS, "pic"))
                blips += len(descendants(root, A_NS, "blip"))
        self.assertEqual((leaves, groups, leaves + groups, pics, blips), (570, 63, 633, 0, 0))

    def test_geometry_only_build_omits_text_and_does_not_create_images(self):
        elements = self._elements("namespace-text.svg", include_text=False)
        self.assertEqual(elements, ())
        nested = self._elements("nested-groups.svg", include_text=False)
        self.assertEqual(len(nested), 1)
        self.assertFalse(descendants(nested[0], A_NS, "t"))
        self.assertFalse(descendants(nested[0], P_NS, "pic"))


class CandidateDeckTests(unittest.TestCase):
    def _modules(self):
        import _ppt_editable.drawingml as drawingml
        import _ppt_editable.notes as notes

        return drawingml, notes

    def _deck_plan(self):
        from _ppt_editable.contract import (
            parse_storyboard,
            resolve_slide_sources,
            validate_completed_run,
        )
        from _ppt_editable.svg_parser import preflight_deck

        directory = tempfile.TemporaryDirectory()
        self.addCleanup(directory.cleanup)
        run = Path(directory.name) / "run"
        shutil.copytree(RUN_FIXTURE_ROOT, run)
        context = validate_completed_run(run)
        storyboard = parse_storyboard(context.storyboard_path)
        sources = resolve_slide_sources(context, storyboard)
        return preflight_deck(
            context,
            sources,
            storyboard,
            "sha256:" + "a" * 64,
        )

    def test_speaker_notes_format_attach_and_extract_semantic_body(self):
        from pptx import Presentation

        _, notes = self._modules()
        from _ppt_editable.contract import StoryboardSlide

        projected = notes.speaker_notes_from_storyboard(
            StoryboardSlide("S01", "结论", "要点", "S02")
        )
        self.assertEqual(
            projected,
            SpeakerNotes("结论", "要点", "S02"),
        )
        value = SpeakerNotes(
            assertion_title="结论",
            audience_takeaway="要点",
            next_link="S02",
        )
        self.assertEqual(
            notes.format_speaker_notes(value),
            "本页结论：结论\n听众要点：要点\n衔接下一页：S02",
        )
        self.assertEqual(
            notes.format_speaker_notes(
                SpeakerNotes(
                    assertion_title="最终结论",
                    audience_takeaway="记住结论",
                    next_link="END",
                )
            ),
            "本页结论：最终结论\n听众要点：记住结论\n衔接下一页：END",
        )
        self.assertEqual(
            notes.format_speaker_notes(SpeakerNotes(assertion_title="仅结论")),
            "本页结论：仅结论",
        )
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        notes.attach_speaker_notes(slide, value)
        path = Path(tempfile.mkdtemp()) / "notes.pptx"
        self.addCleanup(lambda: shutil.rmtree(path.parent, ignore_errors=True))
        presentation.save(path)
        self.assertEqual(
            notes.extract_speaker_notes(path, 0),
            ("本页结论：结论", "听众要点：要点", "衔接下一页：S02"),
        )

    def test_missing_optional_note_field_is_recorded_as_deck_warning(self):
        from _ppt_editable.contract import (
            parse_storyboard,
            resolve_slide_sources,
            validate_completed_run,
        )
        from _ppt_editable.svg_parser import preflight_deck

        directory = tempfile.TemporaryDirectory()
        self.addCleanup(directory.cleanup)
        run = Path(directory.name) / "run"
        shutil.copytree(RUN_FIXTURE_ROOT, run)
        storyboard_path = run / ".ppt-pilot" / "故事板.md"
        storyboard_path.write_text(
            storyboard_path.read_text(encoding="utf-8").replace(
                "- **next_link**: S02\n",
                "",
            ),
            encoding="utf-8",
        )
        context = validate_completed_run(run)
        storyboard = parse_storyboard(context.storyboard_path)
        sources = resolve_slide_sources(context, storyboard)
        deck = preflight_deck(
            context,
            sources,
            storyboard,
            "sha256:" + "a" * 64,
        )
        self.assertIn("S01: missing optional note field next_link", deck.warnings)

    def test_build_presentation_uses_exact_dimensions_titles_shapes_and_notes(self):
        drawingml, notes = self._modules()
        deck = self._deck_plan()
        presentation = drawingml.build_presentation(deck)
        self.assertEqual((presentation.slide_width, presentation.slide_height), (12192000, 6858000))
        self.assertEqual(len(presentation.slides), 2)
        for index, (slide, plan) in enumerate(zip(presentation.slides, deck.slides)):
            self.assertEqual(slide._element.cSld.get("name"), plan.title)
            self.assertGreater(len(slide.shapes), 0)
            self.assertFalse(any(shape.is_placeholder for shape in slide.shapes))
            expected_notes = tuple(
                line
                for line in notes.format_speaker_notes(plan.notes).splitlines()
                if line
            )
            actual_notes = tuple(
                paragraph.text
                for paragraph in slide.notes_slide.notes_text_frame.paragraphs
                if paragraph.text
            )
            self.assertEqual(actual_notes, expected_notes, index)
            ids = [
                int(node.get("id"))
                for node in slide._element.findall(".//" + qn(P_NS, "cNvPr"))
            ]
            self.assertEqual(len(ids), len(set(ids)))

    def test_presentation_bytes_are_zip_valid_reopenable_and_image_free(self):
        from pptx import Presentation

        drawingml, notes = self._modules()
        deck = self._deck_plan()
        data = drawingml.presentation_bytes(deck)
        self.assertTrue(data.startswith(b"PK"))
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            self.assertIsNone(archive.testzip())
            names = archive.namelist()
            self.assertFalse(any(name.startswith("ppt/media/") for name in names))
            self.assertEqual(len([name for name in names if re.fullmatch(r"ppt/slides/slide[0-9]+\.xml", name)]), 2)
            self.assertEqual(len([name for name in names if re.fullmatch(r"ppt/notesSlides/notesSlide[0-9]+\.xml", name)]), 2)
        reopened = Presentation(io.BytesIO(data))
        self.assertEqual(len(reopened.slides), len(deck.slides))
        self.assertEqual((reopened.slide_width, reopened.slide_height), (12192000, 6858000))
        for index, (slide, plan) in enumerate(zip(reopened.slides, deck.slides)):
            self.assertEqual(slide._element.cSld.get("name"), plan.title, index)
            expected_notes = tuple(
                line
                for line in notes.format_speaker_notes(plan.notes).splitlines()
                if line
            )
            actual_notes = tuple(
                paragraph.text
                for paragraph in slide.notes_slide.notes_text_frame.paragraphs
                if paragraph.text
            )
            self.assertEqual(actual_notes, expected_notes, index)
            self.assertEqual(
                len(slide.shapes),
                len(drawingml.build_slide(plan, ShapeIdAllocator())),
                index,
            )
        self.assertEqual(
            tuple(
                paragraph.text
                for paragraph in reopened.slides[-1].notes_slide.notes_text_frame.paragraphs
                if paragraph.text
            )[-1],
            "衔接下一页：END",
        )

    @unittest.skipUnless(os.environ.get("PPT_EDITABLE_REFERENCE_RUN"), "reference run not configured")
    def test_reference_candidate_has_14_reopenable_noted_image_free_slides(self):
        from pptx import Presentation
        from _ppt_editable.contract import (
            parse_storyboard,
            resolve_slide_sources,
            validate_completed_run,
        )
        from _ppt_editable.svg_parser import preflight_deck
        import _ppt_editable.drawingml as drawingml

        run = Path(os.environ["PPT_EDITABLE_REFERENCE_RUN"])
        context = validate_completed_run(run)
        storyboard = parse_storyboard(context.storyboard_path)
        sources = resolve_slide_sources(context, storyboard)
        deck = preflight_deck(
            context,
            sources,
            storyboard,
            "sha256:" + "a" * 64,
        )
        data = drawingml.presentation_bytes(deck)
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            self.assertIsNone(archive.testzip())
            self.assertFalse(any(name.startswith("ppt/media/") for name in archive.namelist()))
        reopened = Presentation(io.BytesIO(data))
        self.assertEqual(len(reopened.slides), 14)
        self.assertEqual((reopened.slide_width, reopened.slide_height), (12192000, 6858000))
        for slide in reopened.slides:
            self.assertFalse(any(shape.is_placeholder for shape in slide.shapes))
            self.assertTrue(
                any(
                    paragraph.text
                    for paragraph in slide.notes_slide.notes_text_frame.paragraphs
                )
            )

    def test_duplicate_slide_owner_is_rejected_before_candidate_creation(self):
        drawingml, _ = self._modules()
        deck = self._deck_plan()
        duplicate = DeckPlan(
            deck_id=deck.deck_id,
            input_snapshot_id=deck.input_snapshot_id,
            slides=(deck.slides[0], deck.slides[0]),
            warnings=deck.warnings,
        )
        with self.assertRaisesRegex(ValueError, "duplicate slide owner"):
            drawingml.build_presentation(duplicate)

    def test_geometry_only_serialization_does_not_mutate_deck_plan(self):
        drawingml, _ = self._modules()
        deck = self._deck_plan()
        before = repr(deck)
        data = drawingml.presentation_bytes(deck, include_text=False)
        self.assertTrue(data.startswith(b"PK"))
        self.assertEqual(repr(deck), before)

    def test_presentation_bytes_closes_its_temporary_stream(self):
        drawingml, _ = self._modules()
        deck = self._deck_plan()
        self.assertTrue(
            hasattr(drawingml, "_presentation_stream"),
            "presentation_bytes needs a close-observable stream seam",
        )
        stream = io.BytesIO()
        with mock.patch.object(
            drawingml,
            "_presentation_stream",
            return_value=stream,
        ):
            data = drawingml.presentation_bytes(deck)
        self.assertTrue(data.startswith(b"PK"))
        self.assertTrue(stream.closed)

    def test_geometry_only_presentation_has_no_text_or_empty_top_level_groups(self):
        drawingml, _ = self._modules()
        deck = self._deck_plan()
        presentation = drawingml.build_presentation(deck, include_text=False)
        for slide in presentation.slides:
            self.assertFalse(slide._element.findall(".//" + qn(A_NS, "t")))
            for group in slide._element.findall(".//" + qn(P_NS, "grpSp")):
                children = [
                    child
                    for child in group
                    if child.tag in (qn(P_NS, "sp"), qn(P_NS, "grpSp"))
                ]
                self.assertTrue(children)

    def test_repeated_build_does_not_mutate_deck_plan(self):
        drawingml, _ = self._modules()
        deck = self._deck_plan()
        before = repr(deck)
        first = drawingml.presentation_bytes(deck)
        second = drawingml.presentation_bytes(deck)
        self.assertEqual(repr(deck), before)
        self.assertEqual(first, second)
        self.assertTrue(first.startswith(b"PK"))


if __name__ == "__main__":
    unittest.main()
